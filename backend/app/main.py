# backend/app/main.py
from typing import List
import os
import shutil

from fastapi import FastAPI, Body, HTTPException, Query
from fastapi.middleware.cors import CORSMiddleware
from fastapi_limiter import FastAPILimiter
from fastapi_limiter.depends import RateLimiter
from openai import OpenAI


from pptx import Presentation
from pptx.util import Inches, Pt
from datetime import datetime
from io import BytesIO
from fastapi.responses import StreamingResponse
from reportlab.lib.pagesizes import A4
from reportlab.lib.units import inch
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak



# our modules
import app.rag_store_chromadb as rag
from app.ingest_dir import ingest_dir
from .models import (
    ChatRequest, ChatResponse, Source, Message,
    GenerateRequest, GenerateResponse
)
from . import deps               # env + config
from .api_router import api_router  # API router with proper structure
from .cleanup_service import cleanup_service

app = FastAPI(title="ICLeaF Chatbot", version="0.2")

# ---- CORS (for the React frontend) ----
# IMPORTANT: Add CORS middleware BEFORE including routers
# CORS middleware automatically handles OPTIONS preflight requests
app.add_middleware(
    CORSMiddleware,
    allow_origins=deps.ALLOWED_ORIGINS if deps.ALLOWED_ORIGINS else ["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Include the API router
app.include_router(api_router, prefix="/api")

# ---- preload your repo docs on boot (Internal mode data) ----
@app.on_event("startup")
async def _preload_docs():
    # Validate OpenAI API key on startup
    if not deps.OPENAI_API_KEY:
        print("[startup] ⚠️  WARNING: OpenAI API key is not set!")
        print("[startup] ⚠️  The application will start but chatbot and content generation features will not work.")
        print("[startup] ⚠️  Please set OPENAI_API_KEY in your .env file and restart the server.")
    else:
        # Test API key by making a simple request
        try:
            from openai import OpenAI
            test_client = OpenAI(api_key=deps.OPENAI_API_KEY)
            # Make a minimal test request
            test_client.models.list()
            print("[startup] ✓ OpenAI API key validated successfully")
        except Exception as e:
            error_msg = str(e)
            if "Invalid API key" in error_msg or "incorrect API key" in error_msg.lower():
                print("[startup] ❌ ERROR: Invalid OpenAI API key!")
                print(f"[startup] ❌ Error: {error_msg}")
                print("[startup] ❌ Please check your .env file and ensure the API key is correct.")
                print("[startup] ❌ The application will start but API calls will fail.")
            elif "Insufficient quota" in error_msg or "quota" in error_msg.lower():
                print("[startup] ⚠️  WARNING: OpenAI API quota issue detected")
                print(f"[startup] ⚠️  Error: {error_msg}")
                print("[startup] ⚠️  Please check your OpenAI account billing.")
            else:
                print(f"[startup] ⚠️  WARNING: Could not validate API key: {error_msg}")
                print("[startup] ⚠️  The application will start but API calls may fail.")
    
    # Initialize rate limiter
    try:
        # await FastAPILimiter.init()
        print("[startup] Rate limiter disabled for development")
    except Exception as e:
        print(f"[startup] Rate limiter initialization failed: {e}")
        print("[startup] Continuing without rate limiting...")
    
    # Start cleanup service
    cleanup_service.start()
    print("[startup] Cleanup service started")
    
    docs_dir = os.getenv("DOCS_DIR", "./seed_docs")
    reindex = os.getenv("REINDEX_ON_START", "false").lower() == "true"
    skip_ingestion = os.getenv("SKIP_INGESTION_ON_START", "false").lower() == "true"

    if skip_ingestion:
        print(f"[startup] Skipping document ingestion (SKIP_INGESTION_ON_START=true)")
        print(f"[startup] ChromaDB already has {rag.count()} documents")
    elif not os.path.isdir(docs_dir):
        print(f"[startup] DOCS_DIR not found: {docs_dir}, skipping ingestion")
        print(f"[startup] ChromaDB already has {rag.count()} documents")
    else:
        if reindex:
            # delete on-disk index and reset collection
            shutil.rmtree("./data/chroma", ignore_errors=True)
            rag.reset_index()

        print(f"[startup] Starting document ingestion from {docs_dir}...")
        try:
            count_added = ingest_dir(docs_dir)
            print(f"[startup] ✓ Ingested {count_added} chunks from {docs_dir} (REINDEX_ON_START={reindex})")
        except Exception as e:
            print(f"[startup] ⚠️  Error during document ingestion: {e}")
            print(f"[startup] ⚠️  Continuing startup with existing {rag.count()} documents in ChromaDB")
    
    print(f"[startup] ✓ Application startup complete")

@app.on_event("shutdown")
async def _shutdown():
    """Cleanup on shutdown."""
    cleanup_service.stop()
    print("[shutdown] Cleanup service stopped")




def _build_pdf_bytes(
    title: str,
    content: str,
    sources: List[Source],
    meta: dict,
) -> bytes:
    """
    Build a simple, clean PDF from generated content + sources using reportlab.
    """
    buffer = BytesIO()
    doc = SimpleDocTemplate(
        buffer,
        pagesize=A4,
        leftMargin=0.8*inch,
        rightMargin=0.8*inch,
        topMargin=0.8*inch,
        bottomMargin=0.8*inch,
        title=title,
        author="ICLeaF LMS",
    )

    styles = getSampleStyleSheet()
    h1 = styles["Heading1"]
    h2 = styles["Heading2"]
    body = styles["BodyText"]
    mono = ParagraphStyle("Mono", parent=body, fontName="Courier", fontSize=9, leading=11)

    story = []

    # Title
    story.append(Paragraph(title, h1))
    story.append(Spacer(1, 0.2*inch))

    # Metadata block
    when = meta.get("generated_at") or datetime.now().strftime("%Y-%m-%d %H:%M")
    role = meta.get("role", "Learner")
    mode = meta.get("mode", "internal").title()
    kind = meta.get("kind", "summary").title()

    story.append(Paragraph(f"<b>Type:</b> {kind}", body))
    story.append(Paragraph(f"<b>Role:</b> {role}", body))
    story.append(Paragraph(f"<b>Mode:</b> {mode}", body))
    story.append(Paragraph(f"<b>Generated:</b> {when}", body))
    story.append(Spacer(1, 0.25*inch))

    # Content
    story.append(Paragraph("Content", h2))
    story.append(Spacer(1, 0.1*inch))

    # Split on double newlines to keep paragraphs readable
    for para in content.split("\n\n"):
        story.append(Paragraph(para.strip().replace("\n", "<br/>"), body))
        story.append(Spacer(1, 0.08*inch))

    # Sources
    if sources:
        story.append(Spacer(1, 0.25*inch))
        story.append(Paragraph("Sources", h2))
        story.append(Spacer(1, 0.1*inch))
        for i, s in enumerate(sources, 1):
            title = s.title or "Document"
            if s.url:
                # show url in a monospaced smaller line (not clickable in pure text PDF)
                story.append(Paragraph(f"[{i}] {title}", body))
                story.append(Paragraph(str(s.url), mono))
            else:
                story.append(Paragraph(f"[{i}] {title}", body))
            story.append(Spacer(1, 0.04*inch))

    doc.build(story)
    pdf_bytes = buffer.getvalue()
    buffer.close()
    return pdf_bytes






def _build_pptx_bytes(
    title: str,
    content: str,
    sources: List[Source],
    meta: dict,
) -> bytes:
    """
    Build a simple PPTX:
      - Title slide
      - Content broken into slides by heuristics
      - A Sources slide at the end
    """
    prs = Presentation()
    # Common sizes
    title_layout = prs.slide_layouts[0]   # Title
    bullet_layout = prs.slide_layouts[1]  # Title + Content

    # Title slide
    slide = prs.slides.add_slide(title_layout)
    slide.shapes.title.text = title
    subtitle = slide.placeholders[1]
    subtitle.text = f"{meta.get('kind','Summary').title()} • {meta.get('role','Learner')} • {meta.get('mode','internal').title()}"

    # Prepare content blocks
    # Split on double newlines for paragraphs/sections
    blocks = [b.strip() for b in content.split("\n\n") if b.strip()]

    def add_bullet_slide(title_text: str, bullets: list[str]):
        slide = prs.slides.add_slide(bullet_layout)
        slide.shapes.title.text = title_text
        body = slide.shapes.placeholders[1].text_frame
        body.clear()
        for i, line in enumerate(bullets):
            # strip leading list characters
            txt = line.strip().lstrip("-•*").strip()
            if i == 0:
                body.text = txt
            else:
                p = body.add_paragraph()
                p.text = txt
                p.level = 0
        # make text a bit larger
        for p in body.paragraphs:
            for run in p.runs:
                run.font.size = Pt(18)

    # Heuristic: if a block is a bullet list (multiple lines starting with -/*/•), keep them on one slide.
    current_bullets = []
    current_title = "Content"
    slide_count = 0

    def flush_bullets():
        nonlocal current_bullets, current_title, slide_count
        if current_bullets:
            add_bullet_slide(current_title if slide_count == 0 else f"{current_title} (cont.)", current_bullets[:10])
            current_bullets = []
            slide_count += 1

    for block in blocks:
        lines = block.split("\n")
        is_bullety = sum(1 for ln in lines if ln.strip().startswith(("-", "*", "•"))) >= max(2, len(lines)//2)

        if is_bullety:
            # add as one bullets slide
            flush_bullets()
            bullets = [ln for ln in lines if ln.strip()]
            add_bullet_slide("Key Points", bullets[:10])
            slide_count += 1
        else:
            # normal paragraph; collect and wrap bullets every ~8 lines
            for ln in lines:
                if not ln.strip():
                    continue
                # break long paragraphs into pseudo bullets
                current_bullets.append(ln.strip())
                if len(current_bullets) >= 8:
                    flush_bullets()
    flush_bullets()

    # Sources slide
    if sources:
        slide = prs.slides.add_slide(bullet_layout)
        slide.shapes.title.text = "Sources"
        body = slide.shapes.placeholders[1].text_frame
        body.clear()
        for i, s in enumerate(sources, 1):
            title_line = f"[{i}] {s.title or 'Document'}"
            if i == 1:
                body.text = title_line
            else:
                p = body.add_paragraph()
                p.text = title_line
                p.level = 0
            # add URL as a sub-bullet if present
            if s.url:
                p2 = body.add_paragraph()
                p2.text = str(s.url)
                p2.level = 1
        for p in body.paragraphs:
            for run in p.runs:
                run.font.size = Pt(16)

    bio = BytesIO()
    prs.save(bio)
    return bio.getvalue()



# ========== UVICORN ENTRY POINT ==========
# Required for running: python backend/app/main.py

if __name__ == "__main__":
    import uvicorn
    
    # Get configuration from environment or use defaults
    host = os.getenv("BACKEND_HOST", "0.0.0.0")  # Bind to all interfaces
    port = int(os.getenv("BACKEND_PORT", 8000))
    reload = os.getenv("RELOAD", "true").lower() == "true"
    
    print("\n" + "="*70)
    print("🚀 Starting ICLeaF LMS AI Backend Server")
    print("="*70)
    print(f"📍 Host: {host}")
    print(f"📍 Port: {port}")
    print(f"🔄 Auto-reload: {reload}")
    print(f"📚 Docs: http://localhost:{port}/docs")
    print("="*70 + "\n")
    
    uvicorn.run(
        "app.main:app",
        host=host,
        port=port,
        reload=reload,
        log_level="info"
    )












# CORS middleware is now configured at the top (before router inclusion)

# ---- LLM client ----
client = OpenAI(api_key=deps.OPENAI_API_KEY) if deps.OPENAI_API_KEY else None

# Health and stats endpoints are now handled by api_router

@app.get("/internal/search")
def internal_search(q: str = Query(..., min_length=2), top_k: int = 6):
    hits = rag.query(q, top_k=top_k)
    results = []
    for h in hits:
        meta = h.get("meta", {})
        results.append({
            "snippet": h["text"][:400],
            "title": meta.get("title", meta.get("filename", "Document")),
            "filename": meta.get("filename"),
            "page": meta.get("page"),
            "score": h.get("score"),
        })
    return {"ok": True, "q": q, "results": results}

@app.get("/internal/topics")
def internal_topics(top_n: int = 20):
    try:
        from sklearn.feature_extraction.text import TfidfVectorizer
    except Exception as e:
        return {"ok": False, "msg": f"scikit-learn not installed: {e}"}

    docs, metas = rag.all_documents()
    if not docs:
        return {"ok": True, "topics": []}

    vec = TfidfVectorizer(max_features=5000, ngram_range=(1,2), stop_words="english")
    X = vec.fit_transform(docs)
    means = X.mean(axis=0).A1
    feats = vec.get_feature_names_out()
    top = sorted(zip(means, feats), reverse=True)[:top_n]
    topics = [t for _, t in top]
    return {"ok": True, "topics": topics}

@app.post("/reindex")
def reindex(
    docs_dir: str = Body(None),
    drop_db: bool = Body(False),
):
    """
    Force reindex:
      - optionally remove ./data/chroma on disk (drop_db=True)
      - reset the Chroma collection
      - ingest docs from docs_dir (or DOCS_DIR env)
    """
    root = docs_dir or os.getenv("DOCS_DIR", "./seed_docs")
    if drop_db:
        shutil.rmtree("./data/chroma", ignore_errors=True)
    rag.reset_index()

    if not os.path.isdir(root):
        return {"ok": False, "msg": f"docs_dir not found: {root}", "index_count": rag.count()}

    added = ingest_dir(root)
    return {"ok": True, "docs_dir": root, "chunks_added": added, "index_count": rag.count()}

# ========= CONTENT GENERATION =========
@app.post("/generate", response_model=GenerateResponse)
async def generate(req: GenerateRequest = Body(...)):
    """
    Generate educational content (summary, quiz, or lesson plan)
    using either Internal (RAG) or Cloud context.
    """
    if client is None:
        raise HTTPException(status_code=500, detail="OPENAI_API_KEY not configured")

    # --------- Gather context + sources ----------
    context_blocks: List[str] = []
    sources: List[Source] = []

    if req.mode == "internal":
        hits = rag.query(req.topic, top_k=req.top_k)
        for h in hits:
            meta = h.get("meta", {})
            sources.append(
                Source(
                    title=meta.get("title", meta.get("filename", "Document")),
                    url=meta.get("url"),
                    score=h.get("score"),
                )
            )
            context_blocks.append(h["text"])
    else:  # cloud
        if deps.TAVILY_API_KEY:
            try:
                web_results = await wc.tavily_search(req.topic, deps.TAVILY_API_KEY, max_results=req.max_context_blocks)
                for r in web_results:
                    sources.append(Source(title=r.get("title") or r.get("url", "Web page"),
                                          url=r.get("url"), score=r.get("score")))
                    if r.get("url"):
                        try:
                            txt = await wc.fetch_url_text(r["url"])
                            if txt:
                                context_blocks.append(txt)
                        except Exception:
                            continue
            except Exception:
                pass

    if len(context_blocks) > req.max_context_blocks:
        context_blocks = context_blocks[:req.max_context_blocks]

    # --------- Build prompts ----------
    role_hint = {
        "Learner":  "Use beginner-friendly language, short paragraphs, and examples.",
        "Trainer":  "Be detailed, structured, and include objectives/outcomes.",
        "Admin":    "Be concise, structured, and highlight compliance/policy relevance.",
    }.get(req.role, "Use clear and concise language.")

    if req.kind == "summary":
        task_hint = (
            "Write a crisp 6–10 bullet summary of the key concepts. "
            "Avoid hallucinations; only use the provided context. "
            "Finish with 2–3 practical tips or cautions."
        )
    elif req.kind == "quiz":
        task_hint = (
            f"Generate {req.num_questions} multiple-choice questions (MCQs) with 1 correct answer and 3 plausible distractors. "
            "Vary difficulty. After the list, include an **Answer Key** mapping Q# to the correct option. "
            "Only use facts present in the context."
        )
    else:  # lesson
        task_hint = (
            "Create a compact lesson plan with sections: Objectives, Prerequisites, Outline (with time boxes), "
            "Hands-on Activity, Assessment, and Further Reading. Only use the provided context."
        )

    system_prompt = (
        "You are ICLeaF LMS's content generator.\n"
        f"Audience role: {req.role}.\n"
        f"{role_hint}\n"
        "If the answer cannot be supported by the context, say you don't know."
    )

    ctx = "\n\n".join([f"[Block {i+1}]\n{b}" for i, b in enumerate(context_blocks)]) if context_blocks else ""
    user_prompt = f"Task: {req.kind} for topic: '{req.topic}'.\n\nFollow these rules:\n- {task_hint}\n- Cite nothing beyond context."

    messages: List[dict] = [{"role": "system", "content": system_prompt}]
    if ctx:
        messages.append({"role": "system", "content": f"Context:\n{ctx}"})
    messages.append({"role": "user", "content": user_prompt})

    completion = client.chat.completions.create(
        model=deps.OPENAI_MODEL,
        messages=messages,
        temperature=0.3 if req.kind != "quiz" else 0.2,
    )
    content = completion.choices[0].message.content

    if not context_blocks:
        content = (
            "I don't have enough context for this topic. "
            f"Try switching mode or reindexing your internal docs. Topic: {req.topic}"
        )

    return GenerateResponse(ok=True, kind=req.kind, topic=req.topic, content=content, sources=sources)





@app.post("/generate/pdf")
async def generate_pdf(req: GenerateRequest = Body(...)):
    """
    Generate content (using the same rules as /generate) and return it as a PDF.
    """
    # Reuse the main generate() to avoid code duplication
    gen = await generate(req)  # returns GenerateResponse
    filename_safe_topic = "".join(c for c in req.topic if c.isalnum() or c in (" ", "_", "-")).strip().replace(" ", "_")
    filename = f"{req.kind}_{filename_safe_topic or 'content'}.pdf"

    meta = {
        "role": req.role,
        "mode": req.mode,
        "kind": req.kind,
        "generated_at": datetime.now().strftime("%Y-%m-%d %H:%M"),
    }
    pdf_bytes = _build_pdf_bytes(
        title=f"{req.kind.title()}: {req.topic}",
        content=gen.content,
        sources=gen.sources,
        meta=meta,
    )

    return StreamingResponse(
        BytesIO(pdf_bytes),
        media_type="application/pdf",
        headers={"Content-Disposition": f'attachment; filename="{filename}"'},
    )






@app.post("/generate/pptx")
async def generate_pptx(req: GenerateRequest = Body(...)):
    """
    Generate content (using /generate rules) and return it as a PPTX deck.
    """
    gen = await generate(req)  # reuse main logic
    filename_safe_topic = "".join(c for c in req.topic if c.isalnum() or c in (" ", "_", "-")).strip().replace(" ", "_")
    filename = f"{req.kind}_{filename_safe_topic or 'content'}.pptx"

    meta = {
        "role": req.role,
        "mode": req.mode,
        "kind": req.kind,
    }
    pptx_bytes = _build_pptx_bytes(
        title=f"{req.kind.title()}: {req.topic}",
        content=gen.content,
        sources=gen.sources,
        meta=meta,
    )

    return StreamingResponse(
        BytesIO(pptx_bytes),
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f'attachment; filename="{filename}"'},
    )







# Chat functionality is now handled by api_router at /api/chatbot/query
