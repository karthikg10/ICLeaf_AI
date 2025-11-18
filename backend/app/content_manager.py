# backend/app/content_manager.py
import os
import uuid
import json
import asyncio
import shutil
import requests
from typing import Dict, List, Optional, Tuple, Union, Any
from datetime import datetime
from reportlab.lib.pagesizes import letter, A4
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.lib import colors
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import csv
import xlsxwriter
import re
from .models import (
    GeneratedContent, GenerateContentRequest,
    ContentType, ContentStatus, FlashcardConfig, QuizConfig,
    AssessmentConfig, VideoConfig, AudioConfig, CompilerConfig,
    PDFConfig, PPTConfig
)
from . import rag_store_chromadb as rag
from . import web_cloud as wc
from . import deps
from . import media_generator
from openai import OpenAI

# In-memory content storage (in production, use database)
_content_storage: Dict[str, GeneratedContent] = {}

# Download tracking (in production, use database)
_download_tracking: Dict[str, List[Dict]] = {}  # contentId -> [{"timestamp": ..., "userId": ..., "ip": ...}]

# Content generation client
content_client = OpenAI(api_key=deps.OPENAI_API_KEY) if deps.OPENAI_API_KEY else None


def extract_openai_response(response):
    """Safely extract content from OpenAI API response with full debugging."""
    try:
        # FIRST: Check if response has error
        if hasattr(response, 'error') and response.error is not None:
            error_detail = response.error
            print(f"[ERROR] OpenAI returned error: {error_detail}")
            raise ValueError(f"OpenAI API Error: {error_detail}")
        
        # SECOND: Check choices exist
        if not hasattr(response, 'choices'):
            print(f"[ERROR] Response has no 'choices' attribute")
            print(f"[DEBUG] Response attributes: {list(vars(response).keys())}")
            raise ValueError("Response missing 'choices'")
        
        if len(response.choices) == 0:
            print(f"[ERROR] Response.choices is empty")
            raise ValueError("No choices in response")
        
        print(f"[DEBUG] ✓ Response has {len(response.choices)} choices")
        
        # THIRD: Get first choice
        choice = response.choices[0]
        
        if not hasattr(choice, 'message'):
            print(f"[ERROR] Choice has no 'message' attribute")
            print(f"[DEBUG] Choice type: {type(choice)}")
            print(f"[DEBUG] Choice attributes: {list(vars(choice).keys())}")
            raise ValueError("Choice missing 'message'")
        
        print(f"[DEBUG] ✓ Choice has message")
        
        # FOURTH: Get message
        message = choice.message
        
        if not hasattr(message, 'content'):
            print(f"[ERROR] Message has no 'content' attribute")
            print(f"[DEBUG] Message type: {type(message)}")
            print(f"[DEBUG] Message attributes: {list(vars(message).keys())}")
            raise ValueError("Message missing 'content'")
        
        print(f"[DEBUG] ✓ Message has content")
        
        # FIFTH: Get content
        content = message.content
        
        if not content:
            print(f"[ERROR] Content is empty or None")
            raise ValueError("Content is empty")
        
        if not isinstance(content, str):
            print(f"[ERROR] Content is not string, it's {type(content)}")
            content = str(content)
        
        print(f"[DEBUG] ✓ Got content: {len(content)} characters")
        return content
        
    except ValueError as e:
        print(f"[ERROR] ValueError: {str(e)}")
        raise
    except Exception as e:
        print(f"[ERROR] Unexpected error: {type(e).__name__}: {str(e)}")
        raise ValueError(f"Failed to extract response: {str(e)}")


def validate_custom_path(custom_path: str) -> bool:
    """Validate custom file path for security."""
    if not custom_path:
        return True
    
    # Prevent directory traversal
    if ".." in custom_path:
        raise ValueError("Invalid path: directory traversal (..) not allowed")
    
    if custom_path.startswith("/etc") or custom_path.startswith("/sys"):
        raise ValueError("Invalid path: system directories not allowed")
    
    return True


def generate_content_id() -> str:
    """Generate a unique content ID."""
    return str(uuid.uuid4())

def clean_markdown_formatting(text: str) -> str:
    """Remove ALL markdown formatting from text - clean plain text only."""
    
    # Remove ** and __ (bold)
    text = re.sub(r'\*\*(.*?)\*\*', r'\1', text)
    text = re.sub(r'__(.*?)__', r'\1', text)
    
    # Remove * and _ (italic) - but be careful with contractions
    text = re.sub(r'([^a-zA-Z0-9])\*(.*?)\*([^a-zA-Z0-9])', r'\1\2\3', text)
    text = re.sub(r'([^a-zA-Z0-9])_(.*?)_([^a-zA-Z0-9])', r'\1\2\3', text)
    
    # Remove ` (inline code)
    text = re.sub(r'`(.*?)`', r'\1', text)
    
    # Remove ``` (code blocks)
    text = re.sub(r'```.*?```', '', text, flags=re.DOTALL)
    
    # Remove # (headers)
    text = re.sub(r'^#+\s+', '', text, flags=re.MULTILINE)
    
    # Remove - (lists) - replace with just text
    text = re.sub(r'^\s*[-*+]\s+', '', text, flags=re.MULTILINE)
    
    # Remove > (blockquotes)
    text = re.sub(r'^>\s+', '', text, flags=re.MULTILINE)
    
    # Remove [text](url) links - keep just text
    text = re.sub(r'\[([^\]]+)\]\([^)]+\)', r'\1', text)
    
    # Remove --- (horizontal rules)
    text = re.sub(r'^---+$', '', text, flags=re.MULTILINE)
    
    # Remove | (tables)
    text = re.sub(r'\|.*\|', '', text)
    
    # Clean up extra whitespace
    text = re.sub(r'\n\s*\n\s*\n', '\n\n', text)  # Max 2 newlines
    text = re.sub(r' +', ' ', text)  # Single spaces
    text = text.strip()
    
    return text

def get_rag_context_for_internal_mode(request: GenerateContentRequest, top_k: int = 5) -> str:
    """Get RAG context from uploaded documents for internal mode content generation."""
    if request.mode != "internal":
        return ""
    
    try:
        # Query RAG store with the prompt
        # Use -0.2 as min_similarity to be more permissive (same as /api/internal/search)
        # This ensures we get results even when similarity scores are low
        hits = rag.query(
            request.prompt,
            top_k=top_k,
            min_similarity=-0.2
        )
        
        if not hits:
            return ""
        
        # Clean and extract context from hits
        context_blocks = []
        for h in hits:
            text_content = h.get("text", "")
            if not text_content or not isinstance(text_content, str):
                continue
            
            # Clean the text
            text_content = clean_markdown_formatting(text_content)
            text_content = text_content.strip()
            
            if not text_content:
                continue
            
            # Remove excessive whitespace
            text_content = " ".join(text_content.split())
            context_blocks.append(text_content)
        
        if not context_blocks:
            return ""
        
        # Build context string
        ctx_parts = []
        for i, block in enumerate(context_blocks, 1):
            ctx_parts.append(f"[Content Block {i}]\n{block}")
        
        context = "\n\n---\n\n".join(ctx_parts)
        print(f"[CONTENT] Internal mode: Retrieved {len(context_blocks)} context blocks from uploaded documents")
        return context
        
    except Exception as e:
        print(f"[CONTENT] Error getting RAG context: {e}")
        return ""

def get_content_storage_path(user_id: str, content_id: str) -> str:
    """Get the storage path for content."""
    return f"/data/content/{user_id}/{content_id}"

def create_content_directory(user_id: str, content_id: str) -> str:
    """Create directory for content storage."""
    base_path = f"./data/content/{user_id}/{content_id}"
    os.makedirs(base_path, exist_ok=True)
    return base_path

def store_content(content: GeneratedContent) -> None:
    """Store content in memory and optionally on disk."""
    _content_storage[content.contentId] = content

def get_content(content_id: str) -> Optional[GeneratedContent]:
    """Get content by ID."""
    return _content_storage.get(content_id)

def get_user_content(user_id: str, status: Optional[ContentStatus] = None) -> List[GeneratedContent]:
    """Get all content for a user, optionally filtered by status."""
    user_content = [c for c in _content_storage.values() if c.userId == user_id]
    if status:
        user_content = [c for c in user_content if c.status == status]
    return sorted(user_content, key=lambda x: x.createdAt, reverse=True)

def update_content_status(content_id: str, status: ContentStatus, file_path: str = None, error: str = None) -> None:
    """Update content status."""
    if content_id in _content_storage:
        content = _content_storage[content_id]
        content.status = status
        if file_path:
            content.filePath = file_path
            content.downloadUrl = f"/api/content/download/{content_id}"
        if error:
            content.error = error
        if status == "completed":
            content.completedAt = datetime.now()

async def generate_flashcard_content(request: GenerateContentRequest, config: FlashcardConfig) -> str:
    """Generate flashcard content in a clean table format."""
    print(f"[DEBUG] generate_flashcard_content called with config: {config}")
    print(f"[DEBUG] config type: {type(config)}")
    print(f"[DEBUG] config.front: {config.front}")
    
    if not content_client:
        raise Exception("OpenAI client not configured")
    
    # Get RAG context for internal mode
    rag_context = get_rag_context_for_internal_mode(request, top_k=5)
    
    system_prompt = f"""You are an educational content generator. Create flashcards in a clean table format with the following specifications:
- Create 5-8 flashcards based on the topic
- Each flashcard should have a KEY (term/concept) and Description (explanation)
- Format as a clean table with two columns: KEY and Description
- Difficulty: {config.difficulty}
- User role: {request.role}"""
    
    if rag_context:
        system_prompt += f"""

IMPORTANT: You are in INTERNAL MODE. Use the provided context from uploaded documents below to create accurate flashcards. Base your content ONLY on the information provided in the context blocks.

CONTEXT FROM UPLOADED DOCUMENTS:
{rag_context}

Instructions:
- Create flashcards based on the information in the context above
- Use specific terms and concepts found in the documents
- Ensure accuracy to the source material
- If the context doesn't cover the topic fully, create flashcards based on what is available"""
    else:
        system_prompt += """

Format the output as a clean table like this:
| KEY | Description |
|-----|-------------|
| Term 1 | Clear explanation of term 1 |
| Term 2 | Clear explanation of term 2 |

Make sure the table is well-structured and easy to read."""
    
    user_prompt = f"""Create flashcards based on: {request.prompt}
Front: {config.front}
Back: {config.back}
Difficulty: {config.difficulty}

Generate 5-8 flashcards in a clean table format with KEY and Description columns."""
    
    response = content_client.chat.completions.create(
        model=deps.OPENAI_MODEL,
        messages=[
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_prompt}
        ],
        temperature=0.7
    )
    
    return extract_openai_response(response)

def _parse_markdown_table_to_rows(table_text: str) -> List[Tuple[str, str]]:
    """Parse a markdown table with two columns (KEY, Description) into rows."""
    rows: List[Tuple[str, str]] = []
    lines = [ln.strip() for ln in table_text.splitlines() if ln.strip()]
    for ln in lines:
        if '|' not in ln:
            continue
        # Skip markdown separator rows like: |-----|------|
        if set(ln.replace('|', '').replace(' ', '').replace(':', '')) <= {'-',}:
            continue
        parts = [p.strip() for p in ln.split('|')]
        # Typical markdown rows start and end with '|', remove empties created by split
        parts = [p for p in parts if p != '']
        if len(parts) < 2:
            continue
        # Skip a possible header row if it looks like headers
        if parts[0].lower() in ("key", "term") and parts[1].lower().startswith("description"):
            continue
        key_col, desc_col = parts[0], parts[1]
        if key_col and desc_col:
            rows.append((key_col, desc_col))
    return rows

def _write_flashcards_csv_xlsx(storage_path: str, rows: List[Tuple[str, str]]) -> Tuple[str, str]:
    """Write flashcards to CSV and XLSX files WITHOUT header row. Returns (csv_path, xlsx_path)."""
    csv_path = os.path.join(storage_path, "flashcards.csv")
    xlsx_path = os.path.join(storage_path, "flashcards.xlsx")

    # CSV - NO HEADER ROW
    with open(csv_path, 'w', newline='', encoding='utf-8') as f:
        writer = csv.writer(f)
        # Write data directly without header
        for k, d in rows:
            writer.writerow([k, d])

    # XLSX - NO HEADER ROW
    workbook = xlsxwriter.Workbook(xlsx_path)
    worksheet = workbook.add_worksheet("Flashcards")
    
    # Write data directly without header (start from row 0)
    for idx, (k, d) in enumerate(rows):
        worksheet.write(idx, 0, k)
        worksheet.write(idx, 1, d)
    
    # Make columns a bit wider
    worksheet.set_column(0, 0, 28)
    worksheet.set_column(1, 1, 80)
    workbook.close()

    return csv_path, xlsx_path

async def generate_quiz_content(request: GenerateContentRequest, config: QuizConfig) -> str:
    """Generate quiz content in a structured table format."""
    if not content_client:
        raise Exception("OpenAI client not configured")
    
    system_prompt = f"""You are an educational quiz generator. Create a quiz in a structured table format with the following specifications:
- Number of questions: {config.num_questions}
- Difficulty: {config.difficulty}
- Question types: {', '.join(config.question_types)}
- User role: {request.role}

Format the output as a structured table with these columns:
| S.No. | QUESTION | CORRECT ANSWER | ANSWER DESC | ANSWER 1 | ANSWER 2 | ANSWER 3 | ANSWER 4 |

Where:
- S.No.: Serial number (1, 2, 3, etc.)
- QUESTION: The quiz question
- CORRECT ANSWER: The number (1, 2, 3, or 4) indicating which answer is correct
- ANSWER DESC: Brief explanation of why the correct answer is right
- ANSWER 1-4: Four multiple choice options

Make sure each question has exactly 4 answer choices and the correct answer number corresponds to one of them."""
    
    user_prompt = f"""Create a quiz based on: {request.prompt}
Number of questions: {config.num_questions}
Difficulty: {config.difficulty}
Question types: {', '.join(config.question_types)}

Generate the quiz in the exact table format specified above."""
    
    response = content_client.chat.completions.create(
        model=deps.OPENAI_MODEL,
        messages=[
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_prompt}
        ],
        temperature=0.7
    )
    
    return extract_openai_response(response)

async def generate_quiz_table(request: GenerateContentRequest, config: QuizConfig) -> List[Dict[str, Union[str, int]]]:
    """Generate quiz rows suitable for CSV/XLSX export with answer description."""
    if not content_client:
        raise Exception("OpenAI client not configured")

    # Get RAG context for internal mode
    rag_context = get_rag_context_for_internal_mode(request, top_k=5)

    keys = [
        "s_no","question","correct_answer","answer_desc","answer_1","answer_2","answer_3","answer_4"
    ]

    system_prompt = (
        "You generate quiz rows for CSV/XLSX export. Return ONLY JSON array (no markdown) "
        "with objects using EXACT keys: " + ",".join(keys) + ". "
        "'s_no' starts at 1 and increments. 'correct_answer' is 1-4. "
        "'answer_desc' is a brief explanation of why the correct answer is right (20-50 words). "
        f"Create exactly {config.num_questions} rows. Difficulty: {config.difficulty}. "
        "Generate realistic, educational content."
    )
    
    if rag_context:
        system_prompt += f"""

IMPORTANT: You are in INTERNAL MODE. Use the provided context from uploaded documents below to create accurate quiz questions. Base your content ONLY on the information provided in the context blocks.

CONTEXT FROM UPLOADED DOCUMENTS:
{rag_context}

Instructions:
- Create quiz questions based on the information in the context above
- Ensure questions and answers are factually accurate to the source material
- Use specific details from the documents
- If the context doesn't cover the topic fully, create questions based on what is available"""

    user_prompt = (
        f"Create quiz rows for: {request.prompt}. "
        f"Keep questions concise and unambiguous. "
        f"For each row, provide: s_no, question, correct_answer (1-4), "
        f"answer_desc (why the answer is correct), and 4 answer choices. "
        f"Make it realistic and educational."
    )

    response = content_client.chat.completions.create(
        model=deps.OPENAI_MODEL,
        messages=[
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_prompt},
        ],
        temperature=0.3,
    )

    txt = extract_openai_response(response).strip()
    try:
        data = json.loads(txt)
        if not isinstance(data, list):
            raise ValueError("expected list")
    except Exception:
        data = []

    rows: List[Dict[str, Union[str, int]]] = []
    for row in data:
        norm: Dict[str, Union[str, int]] = {}
        for k in keys:
            value = row.get(k, "")
            if k == "s_no":
                # Keep as integer
                try:
                    norm[k] = int(value) if value else 0
                except (ValueError, TypeError):
                    norm[k] = 0
            elif k == "correct_answer":
                # Keep as integer
                try:
                    norm[k] = int(value) if value else 0
                except (ValueError, TypeError):
                    norm[k] = 0
            else:
                # Convert to string
                norm[k] = str(value) if value else ""
        rows.append(norm)
    return rows

def _write_quiz_csv_xlsx(storage_path: str, rows: List[Dict[str, Union[str, int]]]) -> Tuple[str, str]:
    """Write quiz rows to CSV/XLSX with the specified headers order including ANSWER DESC."""
    headers = [
        "S.No.", "QUESTION", "CORRECT ANSWER", "ANSWER DESC",
        "ANSWER 1", "ANSWER 2", "ANSWER 3", "ANSWER 4"
    ]
    key_map = {
        "S.No.": "s_no",
        "QUESTION": "question",
        "CORRECT ANSWER": "correct_answer",
        "ANSWER DESC": "answer_desc",
        "ANSWER 1": "answer_1",
        "ANSWER 2": "answer_2",
        "ANSWER 3": "answer_3",
        "ANSWER 4": "answer_4",
    }

    csv_path = os.path.join(storage_path, "quiz.csv")
    xlsx_path = os.path.join(storage_path, "quiz.xlsx")

    # CSV
    with open(csv_path, 'w', newline='', encoding='utf-8') as f:
        writer = csv.writer(f)
        writer.writerow(headers)
        for r in rows:
            row_data = []
            for h in headers:
                key = key_map[h]
                value = r.get(key, "")
                # Convert S.No. and CORRECT ANSWER to integers for CSV
                if h == "S.No." or h == "CORRECT ANSWER":
                    try:
                        row_data.append(int(value) if value else 0)
                    except (ValueError, TypeError):
                        row_data.append(0)
                else:
                    row_data.append(str(value) if value else "")
            writer.writerow(row_data)

    # XLSX
    workbook = xlsxwriter.Workbook(xlsx_path)
    ws = workbook.add_worksheet("Quiz")
    header_fmt = workbook.add_format({"bold": True})
    for col, h in enumerate(headers):
        ws.write(0, col, h, header_fmt)
    for row_idx, r in enumerate(rows, start=1):
        for col, h in enumerate(headers):
            key = key_map[h]
            value = r.get(key, "")
            # Use write_number for S.No. (column 0) and CORRECT ANSWER (column 2)
            if h == "S.No." or h == "CORRECT ANSWER":
                try:
                    num_value = int(value) if value else 0
                    ws.write_number(row_idx, col, num_value)
                except (ValueError, TypeError):
                    ws.write_number(row_idx, col, 0)
            else:
                ws.write(row_idx, col, str(value) if value else "")
    ws.set_row(0, 18)
    ws.set_column(0, len(headers)-1, 22)
    workbook.close()

    return csv_path, xlsx_path

async def generate_assessment_content(request: GenerateContentRequest, config: AssessmentConfig) -> str:
    """Generate assessment content."""
    if not content_client:
        raise Exception("OpenAI client not configured")
    
    # Generate the assessment content
    system_prompt = f"""You are an educational assessment generator. Create an assessment with the following specifications:
- Duration: {config.duration_minutes} minutes
- Difficulty: {config.difficulty}
- Question types: {', '.join(config.question_types)}
- Passing score: {config.passing_score}%
- User role: {request.role}
- Make it comprehensive and aligned with learning objectives"""
    
    user_prompt = f"""Create an assessment based on: {request.prompt}
Duration: {config.duration_minutes} minutes
Difficulty: {config.difficulty}
Question types: {', '.join(config.question_types)}
Passing score: {config.passing_score}%
Make it comprehensive and aligned with learning objectives."""
    
    response = content_client.chat.completions.create(
        model=deps.OPENAI_MODEL,
        messages=[
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_prompt}
        ],
        temperature=0.7
    )
    
    return extract_openai_response(response)

async def generate_assessment_table(request: GenerateContentRequest, config: AssessmentConfig) -> List[Dict[str, str]]:
    """Generate assessment in exact template format."""
    if not content_client:
        raise Exception("OpenAI client not configured")
    
    # Get RAG context for internal mode
    rag_context = get_rag_context_for_internal_mode(request, top_k=5)
    
    columns = [
        "question", "type", "answer_description", "levels", "total_options",
        "choice_answer_one", "choice_answer_two", "choice_answer_three", 
        "choice_answer_four", "choice_answer_five", "correct_answers", "tag1", "tag2"
    ]
    system_prompt = (
        "You generate assessment rows for CSV/XLSX export. "
        "Return ONLY valid JSON array (no markdown) with EXACTLY these keys: " + ",".join(columns) + ". "
        
        "RULES:\n"
        "1. 'question': Question text (can include #{IMG}, #{VID}, #{FILL})\n"
        "2. 'type': 'Choice', 'FillUp', or 'Match'\n"
        "3. 'answer_description': Brief explanation\n"
        "4. 'levels': 'Easy', 'Medium', or 'Difficult'\n"
        "5. 'total_options': Number of options\n"
        "6. 'choice_answer_one' to 'choice_answer_five': The answer options\n"
        "7. 'correct_answers': '1' or '1,2' for Choice; 'a=1,b=2,c=3' for Match\n"
        "8. 'tag1', 'tag2': Tags\n\n"
        
        f"Create exactly 4 rows. Difficulty: {config.difficulty}."
    )
    
    if rag_context:
        system_prompt += f"""

IMPORTANT: You are in INTERNAL MODE. Use the provided context from uploaded documents below to create accurate assessment questions. Base your content ONLY on the information provided in the context blocks.

CONTEXT FROM UPLOADED DOCUMENTS:
{rag_context}

Instructions:
- Create assessment questions based on the information in the context above
- Ensure questions and answers are factually accurate to the source material
- Use specific details from the documents
- If the context doesn't cover the topic fully, create questions based on what is available"""
    
    user_prompt = (
        f"Create 4 assessment rows for: {request.prompt}\n"
        f"Subject: {request.subjectName or ''}\n"
        f"Topic: {request.topicName or ''}\n"
        f"Duration: {config.duration_minutes} minutes\n"
        f"Difficulty: {config.difficulty}\n"
    )
    response = content_client.chat.completions.create(
        model=deps.OPENAI_MODEL,
        messages=[
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_prompt},
        ],
        temperature=0.4,
    )
    txt = extract_openai_response(response).strip()
    try:
        data = json.loads(txt)
        if not isinstance(data, list):
            raise ValueError("expected list")
    except Exception as e:
        print(f"[ERROR] Failed to parse assessment JSON: {e}")
        data = []
    
    norm_rows: List[Dict[str, str]] = []
    for row in data:
        norm = {}
        for col in columns:
            value = row.get(col, "")
            norm[col] = str(value) if value else ""
        norm_rows.append(norm)
    
    return norm_rows

def _write_assessment_csv_xlsx(storage_path: str, rows: List[Dict[str, str]], subject_name: str = "", topic_name: str = "") -> Tuple[str, str]:
    """Write assessment rows in exact template format."""
    
    csv_path = os.path.join(storage_path, "assessment.csv")
    xlsx_path = os.path.join(storage_path, "assessment.xlsx")

    with open(csv_path, 'w', newline='', encoding='utf-8') as f:
        writer = csv.writer(f)
        
        # Row 1: Subject header only
        row1 = [subject_name] + [''] * 12
        writer.writerow(row1)
        
        # Row 2: Column labels
        row2 = [
            topic_name,
            "type",
            "answer description", 
            "levels",
            "total options",
            '', '', '', '', '', '', '', ''
        ]
        writer.writerow(row2)
        
        # Rows 3+: Data
        for r in rows:
            data_row = [
                r.get("question", ""),
                r.get("type", ""),
                r.get("answer_description", ""),
                r.get("levels", ""),
                r.get("total_options", ""),
                r.get("choice_answer_one", ""),
                r.get("choice_answer_two", ""),
                r.get("choice_answer_three", ""),
                r.get("choice_answer_four", ""),
                r.get("choice_answer_five", ""),
                r.get("correct_answers", ""),
                r.get("tag1", ""),      # ← FIX: Use tag1 from data, NOT subject_name
                r.get("tag2", ""),      # ← FIX: Use tag2 from data, NOT topic_name
            ]
            writer.writerow(data_row)

    # XLSX Format: Same as CSV
    workbook = xlsxwriter.Workbook(xlsx_path)
    ws = workbook.add_worksheet("Sheet1")
    
    # Formats
    subject_fmt = workbook.add_format({
        "bold": True,
        "font_size": 12,
        "bg_color": "#FFFFFF"
    })
    
    label_fmt = workbook.add_format({
        "bold": True,
        "bg_color": "#E7E6E6",
        "border": 1,
        "align": "center",
        "valign": "vcenter"
    })
    
    data_fmt = workbook.add_format({
        "border": 1,
        "align": "left",
        "valign": "top",
        "text_wrap": True
    })
    
    # Row 1: Subject header
    ws.write(0, 0, subject_name, subject_fmt)
    for col in range(1, 13):
        ws.write(0, col, '', subject_fmt)
    
    # Row 2: Column labels
    ws.write(1, 0, topic_name, label_fmt)
    labels = ["type", "answer description", "levels", "total options"] + [''] * 8
    for col, label in enumerate(labels, start=1):
        ws.write(1, col, label, label_fmt)
    
    # Rows 3+: Data
    for row_idx, r in enumerate(rows, start=2):
        data_row = [
            r.get("question", ""),
            r.get("type", ""),
            r.get("answer_description", ""),
            r.get("levels", ""),
            r.get("total_options", ""),
            r.get("choice_answer_one", ""),
            r.get("choice_answer_two", ""),
            r.get("choice_answer_three", ""),
            r.get("choice_answer_four", ""),
            r.get("choice_answer_five", ""),
            r.get("correct_answers", ""),
            r.get("tag1", ""),        # ← FIX: Use tag1 from data
            r.get("tag2", ""),        # ← FIX: Use tag2 from data
        ]
        for col, value in enumerate(data_row):
            ws.write(row_idx, col, value, data_fmt)
    
    # Set column widths
    ws.set_column(0, 0, 20)
    ws.set_column(1, 1, 12)
    ws.set_column(2, 2, 20)
    ws.set_column(3, 3, 10)
    ws.set_column(4, 4, 12)
    ws.set_column(5, 9, 15)
    ws.set_column(10, 10, 12)
    ws.set_column(11, 12, 15)
    
    # Set row heights
    ws.set_row(0, 18)
    ws.set_row(1, 20)
    for i in range(2, len(rows) + 2):
        ws.set_row(i, 35)
    
    workbook.close()

    return csv_path, xlsx_path

async def generate_video_content(request: GenerateContentRequest, config: VideoConfig, content_id: str, storage_path: str) -> str:
    """Generate video script file."""
    if not content_client:
        raise Exception("OpenAI client not configured")
    
    # Get RAG context for internal mode
    rag_context = get_rag_context_for_internal_mode(request, top_k=5)
    
    # Generate the video script
    system_prompt = f"""You are a video content generator. Create engaging video content with the following specifications:
- Duration: {config.duration_seconds} seconds
- Quality: {config.quality}
- Include subtitles: {config.include_subtitles}
- User role: {request.role}
- Make it engaging and suitable for video consumption"""
    
    if rag_context:
        system_prompt += f"""

IMPORTANT: You are in INTERNAL MODE. Use the provided context from uploaded documents below to create accurate video content. Base your content ONLY on the information provided in the context blocks.

CONTEXT FROM UPLOADED DOCUMENTS:
{rag_context}

Instructions:
- Create video content based on the information in the context above
- Ensure the content is factually accurate to the source material
- Use specific details from the documents
- If the context doesn't cover the topic fully, create content based on what is available"""
    
    user_prompt = f"""Create video content based on: {request.prompt}
Duration: {config.duration_seconds} seconds
Quality: {config.quality}
Include subtitles: {config.include_subtitles}
Make it engaging and visually appealing for video viewing."""
    
    response = content_client.chat.completions.create(
        model=deps.OPENAI_MODEL,
        messages=[
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_prompt}
        ],
        temperature=0.7
    )
    
    script_content = extract_openai_response(response)
    
    # Save script file
    try:
        script_path = os.path.join(storage_path, "video_script.txt")
        with open(script_path, 'w', encoding='utf-8') as f:
            f.write(f"Video Script:\n\n{script_content}\n\n")
            f.write(f"Duration: {config.duration_seconds} seconds\n")
            f.write(f"Quality: {config.quality}\n")
            f.write(f"Subtitles: {config.include_subtitles}\n")
            f.write(f"Generated: {datetime.now().isoformat()}\n")
        return script_path
            
    except Exception as e:
        script_path = os.path.join(storage_path, "video_script.txt")
        with open(script_path, 'w', encoding='utf-8') as f:
            f.write(f"Error generating video: {str(e)}\n\nScript: {script_content}")
        return script_path

async def generate_audio_content(request: GenerateContentRequest, config: AudioConfig, content_id: str, storage_path: str) -> str:
    """Generate audio script file and MP3 audio file."""
    if not content_client:
        raise Exception("OpenAI client not configured")
    
    # Get RAG context for internal mode
    rag_context = get_rag_context_for_internal_mode(request, top_k=5)
    
    # Generate the script content
    # IMPORTANT: We need ONLY the spoken script text, no meta-commentary or disclaimers
    system_prompt = f"""You are a professional script writer for audio content. Your task is to write ONLY the spoken script text that will be converted to speech using text-to-speech technology.

CRITICAL RULES:
- Write ONLY the actual spoken words that will be read aloud
- NO meta-commentary, NO disclaimers, NO instructions about audio creation
- NO phrases like "I'm unable to create audio files" or "here's a script"
- Write as if you are the narrator/speaker directly addressing the audience
- The text you write will be directly converted to speech - write it naturally
- Duration target: approximately {config.duration_seconds} seconds when spoken
- Target audience: {request.role}
- Make it engaging, conversational, and suitable for audio listening"""
    
    if rag_context:
        system_prompt += f"""

IMPORTANT: You are in INTERNAL MODE. Use the provided context from uploaded documents below to create an accurate audio script. Base your content ONLY on the information provided in the context blocks.

CONTEXT FROM UPLOADED DOCUMENTS:
{rag_context}

Instructions:
- Create the audio script based on the information in the context above
- Ensure the content is factually accurate to the source material
- Use specific details from the documents
- If the context doesn't cover the topic fully, create content based on what is available"""
    
    user_prompt = f"""Write the spoken script for an audio recording about: {request.prompt}

Requirements:
- Write ONLY the spoken words (no meta-commentary)
- Target duration: {config.duration_seconds} seconds when spoken
- Write naturally as if speaking directly to the listener
- Be engaging and conversational
- Start immediately with the content (no introductions about scripts or audio files)

Begin the script now:"""
    
    response = content_client.chat.completions.create(
        model=deps.OPENAI_MODEL,
        messages=[
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_prompt}
        ],
        temperature=0.7
    )
    
    script_content = extract_openai_response(response)
    original_script = script_content  # Keep original for fallback
    
    # Clean up the script content - remove common meta-commentary patterns
    # Remove disclaimers and meta-text that shouldn't be spoken
    unwanted_patterns = [
        "I'm unable to create audio files",
        "I can provide you with a script",
        "here's a script",
        "here is a script",
        "script that you can use",
        "[Upbeat background music begins]",
        "[Background music",
        "**Narrator:**",
        "**Narrator**",
    ]
    
    # Remove lines that contain unwanted patterns
    lines = script_content.split('\n')
    cleaned_lines = []
    for line in lines:
        line_lower = line.lower().strip()
        # Skip lines that are just separators or contain unwanted patterns
        if any(pattern.lower() in line_lower for pattern in unwanted_patterns):
            continue
        # Skip markdown-style separators (---)
        if line.strip() == '---' or (line.strip().startswith('---') and len(line.strip()) <= 5):
            continue
        # Keep the line if it's actual content
        if line.strip():
            cleaned_lines.append(line)
    
    script_content = '\n'.join(cleaned_lines).strip()
    
    # If the script is empty after cleaning, use the original with minimal cleaning (fallback)
    if not script_content or len(script_content) < 50:
        print("[Audio] Warning: Script content was mostly meta-commentary, using original with minimal cleaning")
        # Try a gentler cleanup - just remove obvious markdown formatting
        script_content = original_script.replace('**', '').replace('---', '').strip()
    
    # Save script file
    try:
        script_path = os.path.join(storage_path, "audio_script.txt")
        with open(script_path, 'w', encoding='utf-8') as f:
            f.write(script_content)
    except Exception as e:
        script_path = os.path.join(storage_path, "audio_script.txt")
        with open(script_path, 'w', encoding='utf-8') as f:
            f.write(f"Error generating audio: {str(e)}\n\nScript: {script_content}")
    
    # Generate MP3 file using TTS
    try:
        # Map voice_type to OpenAI TTS voices
        # Female voices: alloy, nova, shimmer
        # Male voices: echo, fable, onyx
        voice_map = {
            "female": "nova",  # Default female voice
            "male": "onyx"     # Default male voice
        }
        voice = voice_map.get(config.voice_type, "alloy")
        
        # Map quality to TTS model
        # high -> tts-1-hd, medium/low -> tts-1
        model_map = {
            "high": "tts-1-hd",
            "medium": "tts-1",
            "low": "tts-1"
        }
        tts_model = model_map.get(config.quality.lower(), "tts-1")
        
        # Determine output format
        audio_format = config.format.lower() if config.format else "mp3"
        
        # Generate MP3 file path
        audio_path = os.path.join(storage_path, f"audio.{audio_format}")
        
        # Generate the audio file using TTS
        success = await media_generator.generate_audio_file(
            text=script_content,
            output_path=audio_path,
            voice=voice,
            model=tts_model,
            format=audio_format
        )
        
        if success:
            print(f"[Audio] Successfully generated MP3 file: {audio_path}")
            return audio_path
        else:
            print(f"[Audio] Failed to generate MP3, returning script path instead")
            return script_path
            
    except Exception as e:
        print(f"[Audio] Error generating MP3 file: {e}")
        # Return script path as fallback
        return script_path

async def generate_compiler_content(request: GenerateContentRequest, config: CompilerConfig) -> str:
    """Generate compiler/code content."""
    if not content_client:
        raise Exception("OpenAI client not configured")
    
    # Get RAG context for internal mode
    rag_context = get_rag_context_for_internal_mode(request, top_k=5)
    
    system_prompt = f"""You are a code generator. Create code with the following specifications:
- Language: {config.language}
- Include tests: {config.include_tests}
- Difficulty: {config.difficulty}
- User role: {request.role}"""
    
    if rag_context:
        system_prompt += f"""

IMPORTANT: You are in INTERNAL MODE. Use the provided context from uploaded documents below to create accurate code. Base your code ONLY on the information provided in the context blocks.

CONTEXT FROM UPLOADED DOCUMENTS:
{rag_context}

Instructions:
- Create code based on the information in the context above
- Follow coding patterns and examples from the documents
- Ensure the code aligns with the concepts explained in the source material
- If the context doesn't cover the topic fully, create code based on what is available"""
    
    user_prompt = f"""Create code based on: {request.prompt}
Language: {config.language}
Include tests: {config.include_tests}
Difficulty: {config.difficulty}"""
    
    response = content_client.chat.completions.create(
        model=deps.OPENAI_MODEL,
        messages=[
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_prompt}
        ],
        temperature=0.7
    )
    
    return extract_openai_response(response)

async def generate_pdf_content(request: GenerateContentRequest, content_id: str, storage_path: str) -> str:
    """Generate PDF with STRICT page count control."""
    if not content_client:
        raise Exception("OpenAI client not configured")
    
    # Get RAG context for internal mode
    rag_context = get_rag_context_for_internal_mode(request, top_k=5)
    
    pdf_config = request.contentConfig.get('pdf', {})
    num_pages = int(pdf_config.get('num_pages', 5))
    target_audience = pdf_config.get('target_audience', 'general')
    include_images = pdf_config.get('include_images', True)
    difficulty = pdf_config.get('difficulty', 'medium')
    
    print(f"[PDF] Generating EXACTLY {num_pages} pages for {target_audience} at {difficulty} level")
    
    words_per_page = 300
    total_words_needed = num_pages * words_per_page
    
    system_prompt = f"""You are a professional content writer. GENERATE EXACTLY {num_pages} PAGES of content.

STRICT REQUIREMENTS:
1. Write EXACTLY {total_words_needed} words (for {num_pages} pages at ~300 words/page)
2. Target Audience: {target_audience}
3. Difficulty: {difficulty}
4. Structure: Title page + Body sections + Conclusion
5. NO markdown formatting whatsoever
6. NO ** or other special characters
7. Write in clear, proper English sentences
8. Include detailed explanations and examples
9. Break into clear sections with newlines between them

FORMATTING RULES:
- Each section should be 300-400 words
- Use section titles (plain text, no **bold**)
- Separate sections with blank lines
- NO bullet points, NO markdown, NO special formatting
- Write as a professional document

WORD COUNT: {total_words_needed} words minimum"""
    
    if rag_context:
        system_prompt += f"""

IMPORTANT: You are in INTERNAL MODE. Use the provided context from uploaded documents below to create accurate content. Base your content ONLY on the information provided in the context blocks.

CONTEXT FROM UPLOADED DOCUMENTS:
{rag_context}

Instructions:
- Create content based on the information in the context above
- Ensure the content is factually accurate to the source material
- Use specific details and examples from the documents
- Expand on the concepts found in the context to reach the required word count
- If the context doesn't cover the topic fully, create content based on what is available"""

    user_prompt = f"""Create a {num_pages}-page document about: {request.prompt}

MUST REQUIREMENTS:
- Write EXACTLY {total_words_needed} words
- Create {num_pages} full pages of content
- Target audience: {target_audience}
- Difficulty: {difficulty}
- Include detailed explanations
- Include real-world examples
- Professional tone
- NO markdown, NO ** formatting, NO special characters
- Write in plain, clear English

START WRITING NOW - aim for {total_words_needed} words:"""

    try:
        print(f"[PDF] Calling OpenAI API (requesting {total_words_needed} words)...")
        response = content_client.chat.completions.create(
            model=deps.OPENAI_MODEL,
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_prompt}
            ],
            temperature=0.7,
            max_tokens=4000 + (num_pages * 500)
        )
        
        content = extract_openai_response(response)
        print(f"[PDF] Got {len(content)} characters, {len(content.split())} words")
        
    except Exception as e:
        print(f"[PDF] Error: {e}")
        raise
    
    content = clean_markdown_formatting(content)
    content = re.sub(r'\*\*', '', content)
    content = re.sub(r'__', '', content)
    content = re.sub(r'\*(?!\s)', '', content)
    
    pdf_path = os.path.join(storage_path, "document.pdf")
    
    doc = SimpleDocTemplate(
        pdf_path,
        pagesize=A4,
        rightMargin=50,
        leftMargin=50,
        topMargin=50,
        bottomMargin=50
    )
    
    styles = getSampleStyleSheet()
    story = []
    
    heading_style = ParagraphStyle(
        'CustomHeading',
        parent=styles['Heading1'],
        fontSize=18,
        textColor=colors.HexColor('#1a1a4d'),
        spaceAfter=16,
        spaceBefore=12,
        fontName='Helvetica-Bold'
    )
    
    body_style = ParagraphStyle(
        'CustomBody',
        parent=styles['Normal'],
        fontSize=12,
        leading=16,
        spaceAfter=12,
        alignment=4
    )
    
    paragraphs = [p.strip() for p in content.split('\n\n') if p.strip()]
    
    for para in paragraphs:
        if len(para) < 100 and len(para.split()) < 15 and para.isupper():
            story.append(Paragraph(para, heading_style))
            story.append(Spacer(1, 0.1 * inch))
        else:
            story.append(Paragraph(para, body_style))
            story.append(Spacer(1, 0.05 * inch))
    
    try:
        doc.build(story)
        print(f"[PDF] Generated {pdf_path}")
        return pdf_path
    except Exception as e:
        print(f"[PDF] Error building: {e}")
        raise

async def generate_pdf_content_with_path(
    request: GenerateContentRequest, 
    content_id: str, 
    storage_path: str,
    filename: str
) -> str:
    """Generate PDF with custom filename."""
    if not content_client:
        raise Exception("OpenAI client not configured")
    
    # Get RAG context for internal mode
    rag_context = get_rag_context_for_internal_mode(request, top_k=5)
    
    pdf_config = request.contentConfig.get('pdf', {})
    num_pages = int(pdf_config.get('num_pages', 5))
    target_audience = pdf_config.get('target_audience', 'general')
    include_images = pdf_config.get('include_images', True)
    difficulty = pdf_config.get('difficulty', 'medium')
    
    print(f"[PDF] Generating EXACTLY {num_pages} pages for {target_audience} at {difficulty} level")
    print(f"[PDF] Custom filename: {filename}")
    
    words_per_page = 300
    total_words_needed = num_pages * words_per_page
    
    system_prompt = f"""You are a professional content writer. GENERATE EXACTLY {num_pages} PAGES of content.

STRICT REQUIREMENTS:
1. Write EXACTLY {total_words_needed} words (for {num_pages} pages at ~300 words/page)
2. Target Audience: {target_audience}
3. Difficulty: {difficulty}
4. Structure: Title page + Body sections + Conclusion
5. NO markdown formatting whatsoever
6. NO ** or other special characters
7. Write in clear, proper English sentences
8. Include detailed explanations and examples
9. Break into clear sections with newlines between them

FORMATTING RULES:
- Each section should be 300-400 words
- Use section titles (plain text, no **bold**)
- Separate sections with blank lines
- NO bullet points, NO markdown, NO special formatting
- Write as a professional document

WORD COUNT: {total_words_needed} words minimum"""
    
    if rag_context:
        system_prompt += f"""

IMPORTANT: You are in INTERNAL MODE. Use the provided context from uploaded documents below to create accurate content. Base your content ONLY on the information provided in the context blocks.

CONTEXT FROM UPLOADED DOCUMENTS:
{rag_context}

Instructions:
- Create content based on the information in the context above
- Ensure the content is factually accurate to the source material
- Use specific details and examples from the documents
- Expand on the concepts found in the context to reach the required word count
- If the context doesn't cover the topic fully, create content based on what is available"""

    user_prompt = f"""Create a {num_pages}-page document about: {request.prompt}

MUST REQUIREMENTS:
- Write EXACTLY {total_words_needed} words
- Create {num_pages} full pages of content
- Target audience: {target_audience}
- Difficulty: {difficulty}
- Include detailed explanations
- Include real-world examples
- Professional tone
- NO markdown, NO ** formatting, NO special characters
- Write in plain, clear English

START WRITING NOW - aim for {total_words_needed} words:"""

    try:
        print(f"[PDF] Calling OpenAI API (requesting {total_words_needed} words)...")
        response = content_client.chat.completions.create(
            model=deps.OPENAI_MODEL,
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_prompt}
            ],
            temperature=0.7,
            max_tokens=4000 + (num_pages * 500)
        )
        
        content = extract_openai_response(response)
        print(f"[PDF] Got {len(content)} characters, {len(content.split())} words")
        
    except Exception as e:
        print(f"[PDF] Error: {e}")
        raise
    
    content = clean_markdown_formatting(content)
    content = re.sub(r'\*\*', '', content)
    content = re.sub(r'__', '', content)
    
    pdf_path = os.path.join(storage_path, filename)
    
    doc = SimpleDocTemplate(
        pdf_path,
        pagesize=A4,
        rightMargin=50,
        leftMargin=50,
        topMargin=50,
        bottomMargin=50
    )
    
    styles = getSampleStyleSheet()
    story = []
    
    heading_style = ParagraphStyle(
        'CustomHeading',
        parent=styles['Heading1'],
        fontSize=18,
        textColor=colors.HexColor('#1a1a4d'),
        spaceAfter=16,
        spaceBefore=12,
        fontName='Helvetica-Bold'
    )
    
    body_style = ParagraphStyle(
        'CustomBody',
        parent=styles['Normal'],
        fontSize=12,
        leading=16,
        spaceAfter=12,
        alignment=4
    )
    
    paragraphs = [p.strip() for p in content.split('\n\n') if p.strip()]
    
    for para in paragraphs:
        if len(para) < 100 and len(para.split()) < 15 and para.isupper():
            story.append(Paragraph(para, heading_style))
            story.append(Spacer(1, 0.1 * inch))
        else:
            story.append(Paragraph(para, body_style))
            story.append(Spacer(1, 0.05 * inch))
    
    try:
        doc.build(story)
        print(f"[PDF] Generated: {pdf_path}")
        return pdf_path
    except Exception as e:
        print(f"[PDF] Error building: {e}")
        raise

async def generate_ppt_content(request: GenerateContentRequest, content_id: str, storage_path: str) -> str:
    """Generate PowerPoint with ChatGPT-quality rich content."""
    if not content_client:
        raise Exception("OpenAI client not configured")
    
    # Get RAG context for internal mode
    rag_context = get_rag_context_for_internal_mode(request, top_k=5)
    
    ppt_config = request.contentConfig.get('ppt', {})
    num_slides = int(ppt_config.get('num_slides', 10))
    target_audience = ppt_config.get('target_audience', 'general')
    difficulty = ppt_config.get('difficulty', 'medium')
    
    print(f"[PPT] Generating {num_slides} information-rich slides (ChatGPT style)")
    
    system_prompt = f"""You are a world-class presentation expert like ChatGPT. Create {num_slides} RICH, DETAILED slides.

REQUIREMENTS:
- Create EXACTLY {num_slides} slides
- Each slide has a CLEAR TITLE
- Each slide has 4-6 DETAILED bullet points (not just 1-2 words)
- Bullet points are complete sentences with substance
- Include data, examples, explanations
- Professional and engaging
- Suitable for {target_audience} at {difficulty} level

FORMAT FOR EACH SLIDE:
---SLIDE---
TITLE: [Clear, descriptive title]
CONTENT: [4-6 detailed bullet points with substance]
---

SLIDE STRUCTURE:
- Slide 1: Title slide with topic and overview
- Slides 2-{num_slides-1}: Detailed content with examples, data, explanations
- Slide {num_slides}: Summary/Conclusion slide

DETAILED CONTENT RULES:
- Each bullet point should be a complete sentence (20-30 words)
- Include specific examples
- Include relevant statistics or data
- Make it informative, not just pretty
- Rich with actual information"""
    
    if rag_context:
        system_prompt += f"""

IMPORTANT: You are in INTERNAL MODE. Use the provided context from uploaded documents below to create accurate slides. Base your content ONLY on the information provided in the context blocks.

CONTEXT FROM UPLOADED DOCUMENTS:
{rag_context}

Instructions:
- Create slides based on the information in the context above
- Ensure the content is factually accurate to the source material
- Use specific details, examples, and data from the documents
- If the context doesn't cover the topic fully, create slides based on what is available"""

    user_prompt = f"""Create {num_slides} detailed, information-rich slides about: {request.prompt}

Make it like ChatGPT presentations - FULL OF INFORMATION and DETAILED.
Not generic, not simple - RICH CONTENT.

Target audience: {target_audience}
Difficulty: {difficulty}

Each bullet point should have real substance and information.
Each slide should be valuable and detailed.

Generate all {num_slides} slides with rich, detailed content:"""

    try:
        print(f"[PPT] Calling OpenAI API...")
        response = content_client.chat.completions.create(
            model=deps.OPENAI_MODEL,
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_prompt}
            ],
            temperature=0.8,
            max_tokens=5000
        )
        
        content = extract_openai_response(response)
        print(f"[PPT] Got {len(content)} characters")
        
    except Exception as e:
        print(f"[PPT] Error: {e}")
        raise
    
    ppt_path = os.path.join(storage_path, "presentation.pptx")
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    slides_data = []
    current_slide = {"title": "", "bullets": []}
    
    for line in content.split('\n'):
        line = line.strip()
        
        if line == '---SLIDE---':
            continue
        elif line.startswith('TITLE:'):
            current_slide["title"] = line.replace('TITLE:', '').strip()
        elif line.startswith('CONTENT:') or line.startswith('BULLET:') or line.startswith('-'):
            bullet = line.replace('CONTENT:', '').replace('BULLET:', '').replace('-', '', 1).strip()
            if bullet and len(bullet) > 2:
                current_slide["bullets"].append(bullet[:100])
        elif line == '---':
            if current_slide["title"] or current_slide["bullets"]:
                slides_data.append(current_slide)
            current_slide = {"title": "", "bullets": []}
    
    if current_slide["title"] or current_slide["bullets"]:
        slides_data.append(current_slide)
    
    print(f"[PPT] Parsed {len(slides_data)} slides from content")
    
    while len(slides_data) < num_slides:
        slides_data.append({
            "title": f"Additional Information {len(slides_data)}",
            "bullets": ["Detailed content point 1", "Detailed content point 2", "Detailed content point 3"]
        })
    
    blank_layout = prs.slide_layouts[6]
    
    for i, slide_data in enumerate(slides_data[:num_slides]):
        try:
            slide = prs.slides.add_slide(blank_layout)
            
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
            title_frame = title_box.text_frame
            title_frame.word_wrap = True
            title_frame.text = slide_data["title"] or f"Slide {i + 1}"
            
            for paragraph in title_frame.paragraphs:
                paragraph.font.size = Pt(44)
                paragraph.font.bold = True
                paragraph.font.color.rgb = RGBColor(25, 25, 100)
            
            if slide_data["bullets"]:
                content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5.5))
                content_frame = content_box.text_frame
                content_frame.word_wrap = True
                
                for idx, bullet in enumerate(slide_data["bullets"][:6]):
                    if idx == 0:
                        p = content_frame.paragraphs[0]
                    else:
                        p = content_frame.add_paragraph()
                    
                    p.text = f"• {bullet}"
                    p.level = 0
                    p.font.size = Pt(18)
                    p.font.color.rgb = RGBColor(50, 50, 50)
                    p.space_after = Pt(10)
        
        except Exception as e:
            print(f"[PPT] Error on slide {i + 1}: {e}")
            continue
    
    try:
        prs.save(ppt_path)
        print(f"[PPT] Generated {len(prs.slides)} slides: {ppt_path}")
        return ppt_path
    except Exception as e:
        print(f"[PPT] Error saving: {e}")
        raise

async def generate_ppt_content_with_path(
    request: GenerateContentRequest, 
    content_id: str, 
    storage_path: str,
    filename: str
) -> str:
    """Generate PowerPoint with custom filename."""
    if not content_client:
        raise Exception("OpenAI client not configured")
    
    # Get RAG context for internal mode
    rag_context = get_rag_context_for_internal_mode(request, top_k=5)
    
    ppt_config = request.contentConfig.get('ppt', {})
    num_slides = int(ppt_config.get('num_slides', 10))
    target_audience = ppt_config.get('target_audience', 'general')
    difficulty = ppt_config.get('difficulty', 'medium')
    
    print(f"[PPT] Generating {num_slides} slides for {target_audience} at {difficulty} level")
    print(f"[PPT] Custom filename: {filename}")
    
    system_prompt = f"""You are a world-class presentation expert like ChatGPT. Create {num_slides} RICH, DETAILED slides.

REQUIREMENTS:
- Create EXACTLY {num_slides} slides
- Each slide has a CLEAR TITLE
- Each slide has 4-6 DETAILED bullet points (not just 1-2 words)
- Bullet points are complete sentences with substance
- Include data, examples, explanations
- Professional and engaging
- Suitable for {target_audience} at {difficulty} level

FORMAT FOR EACH SLIDE:
---SLIDE---
TITLE: [Clear, descriptive title]
CONTENT: [4-6 detailed bullet points with substance]
---

SLIDE STRUCTURE:
- Slide 1: Title slide with topic and overview
- Slides 2-{num_slides-1}: Detailed content with examples, data, explanations
- Slide {num_slides}: Summary/Conclusion slide

DETAILED CONTENT RULES:
- Each bullet point should be a complete sentence (20-30 words)
- Include specific examples
- Include relevant statistics or data
- Make it informative, not just pretty
- Rich with actual information"""
    
    if rag_context:
        system_prompt += f"""

IMPORTANT: You are in INTERNAL MODE. Use the provided context from uploaded documents below to create accurate slides. Base your content ONLY on the information provided in the context blocks.

CONTEXT FROM UPLOADED DOCUMENTS:
{rag_context}

Instructions:
- Create slides based on the information in the context above
- Ensure the content is factually accurate to the source material
- Use specific details, examples, and data from the documents
- If the context doesn't cover the topic fully, create slides based on what is available"""

    user_prompt = f"""Create {num_slides} detailed, information-rich slides about: {request.prompt}

Make it like ChatGPT presentations - FULL OF INFORMATION and DETAILED.
Not generic, not simple - RICH CONTENT.

Target audience: {target_audience}
Difficulty: {difficulty}

Each bullet point should have real substance and information.
Each slide should be valuable and detailed.

Generate all {num_slides} slides with rich, detailed content:"""

    try:
        print(f"[PPT] Calling OpenAI API...")
        response = content_client.chat.completions.create(
            model=deps.OPENAI_MODEL,
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_prompt}
            ],
            temperature=0.8,
            max_tokens=5000
        )
        
        content = extract_openai_response(response)
        print(f"[PPT] Got {len(content)} characters")
        
    except Exception as e:
        print(f"[PPT] Error: {e}")
        raise
    
    ppt_path = os.path.join(storage_path, filename)
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    slides_data = []
    current_slide = {"title": "", "bullets": []}
    
    for line in content.split('\n'):
        line = line.strip()
        
        if line == '---SLIDE---':
            continue
        elif line.startswith('TITLE:'):
            current_slide["title"] = line.replace('TITLE:', '').strip()
        elif line.startswith('CONTENT:') or line.startswith('BULLET:') or line.startswith('-'):
            bullet = line.replace('CONTENT:', '').replace('BULLET:', '').replace('-', '', 1).strip()
            if bullet and len(bullet) > 2:
                current_slide["bullets"].append(bullet[:100])
        elif line == '---':
            if current_slide["title"] or current_slide["bullets"]:
                slides_data.append(current_slide)
            current_slide = {"title": "", "bullets": []}
    
    if current_slide["title"] or current_slide["bullets"]:
        slides_data.append(current_slide)
    
    print(f"[PPT] Parsed {len(slides_data)} slides")
    
    while len(slides_data) < num_slides:
        slides_data.append({
            "title": f"Additional Information {len(slides_data)}",
            "bullets": ["Detailed content point 1", "Detailed content point 2", "Detailed content point 3"]
        })
    
    blank_layout = prs.slide_layouts[6]
    
    for i, slide_data in enumerate(slides_data[:num_slides]):
        try:
            slide = prs.slides.add_slide(blank_layout)
            
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
            title_frame = title_box.text_frame
            title_frame.word_wrap = True
            title_frame.text = slide_data["title"] or f"Slide {i + 1}"
            
            for paragraph in title_frame.paragraphs:
                paragraph.font.size = Pt(44)
                paragraph.font.bold = True
                paragraph.font.color.rgb = RGBColor(25, 25, 100)
            
            if slide_data["bullets"]:
                content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5.5))
                content_frame = content_box.text_frame
                content_frame.word_wrap = True
                
                for idx, bullet in enumerate(slide_data["bullets"][:6]):
                    if idx == 0:
                        p = content_frame.paragraphs[0]
                    else:
                        p = content_frame.add_paragraph()
                    
                    p.text = f"• {bullet}"
                    p.level = 0
                    p.font.size = Pt(18)
                    p.font.color.rgb = RGBColor(50, 50, 50)
                    p.space_after = Pt(10)
        
        except Exception as e:
            print(f"[PPT] Error on slide {i + 1}: {e}")
            continue
    
    try:
        prs.save(ppt_path)
        print(f"[PPT] Generated {len(prs.slides)} slides: {ppt_path}")
        return ppt_path
    except Exception as e:
        print(f"[PPT] Error saving: {e}")
        raise

async def process_content_generation(request: GenerateContentRequest) -> GeneratedContent:
    """Process content generation with custom filename and path support."""
    content_id = generate_content_id()
    
    # Validate custom path if provided
    if request.customFilePath:
        try:
            validate_custom_path(request.customFilePath)
        except ValueError as e:
            raise Exception(f"Invalid file path: {str(e)}")
    
    # Create content object
    content = GeneratedContent(
        contentId=content_id,
        userId=request.userId,
        role=request.role,
        mode=request.mode,
        contentType=request.contentType,
        prompt=request.prompt,
        status="pending",
        contentConfig=request.contentConfig,
        metadata={
            "docIds": request.docIds,
            "subjectName": request.subjectName,
            "topicName": request.topicName,
            "customFileName": request.customFileName,
            "customFilePath": request.customFilePath
        }
    )
    
    store_content(content)
    
    try:
        # Determine storage path
        if request.customFilePath:
            storage_path = request.customFilePath
            os.makedirs(storage_path, exist_ok=True)
            print(f"[CONTENT] Using custom path: {storage_path}")
        else:
            storage_path = create_content_directory(request.userId, content_id)
            print(f"[CONTENT] Using default path: {storage_path}")
        
        # Determine filename (base name without extension)
        if request.customFileName:
            base_filename = request.customFileName
        else:
            base_filename = f"content_{content_id[:8]}"
        
        print(f"[CONTENT] Base filename: {base_filename}")
        
        # Generate content based on type
        generated_content = ""
        actual_filename = ""
        
        if request.contentType == "flashcard" and request.contentConfig.get('flashcard'):
            flashcard_config_dict = request.contentConfig.get('flashcard')
            flashcard_config = FlashcardConfig(
                front=flashcard_config_dict.get('front', ''),
                back=flashcard_config_dict.get('back', ''),
                difficulty=flashcard_config_dict.get('difficulty', 'medium')
            )
            generated_content = await generate_flashcard_content(request, flashcard_config)
            rows = _parse_markdown_table_to_rows(generated_content)
            csv_path, xlsx_path = _write_flashcards_csv_xlsx(storage_path, rows)
            actual_filename = os.path.basename(xlsx_path)
            file_path = xlsx_path
        elif request.contentType == "quiz" and request.contentConfig.get('quiz'):
            quiz_config_dict = request.contentConfig.get('quiz')
            quiz_config = QuizConfig(
                num_questions=quiz_config_dict.get('num_questions', 5),
                difficulty=quiz_config_dict.get('difficulty', 'medium'),
                question_types=quiz_config_dict.get('question_types', ['multiple_choice'])
            )
            rows = await generate_quiz_table(request, quiz_config)
            csv_path, xlsx_path = _write_quiz_csv_xlsx(storage_path, rows)
            actual_filename = os.path.basename(xlsx_path)
            file_path = xlsx_path
        elif request.contentType == "assessment" and request.contentConfig.get('assessment'):
            assessment_config_dict = request.contentConfig.get('assessment')
            assessment_config = AssessmentConfig(**assessment_config_dict)
            rows = await generate_assessment_table(request, assessment_config)
            
            # Pass subject and topic names
            csv_path, xlsx_path = _write_assessment_csv_xlsx(
                storage_path, 
                rows,
                subject_name=request.subjectName or "Subject",
                topic_name=request.topicName or "Topic"
            )
            
            actual_filename = os.path.basename(xlsx_path)
            file_path = xlsx_path
        elif request.contentType == "video" and request.contentConfig.get('video'):
            video_config_dict = request.contentConfig.get('video')
            video_config = VideoConfig(**video_config_dict)
            generated_content = await generate_video_content(request, video_config, content_id, storage_path)
            actual_filename = os.path.basename(generated_content)
            file_path = generated_content
        elif request.contentType == "audio" and request.contentConfig.get('audio'):
            audio_config_dict = request.contentConfig.get('audio')
            audio_config = AudioConfig(**audio_config_dict)
            generated_content = await generate_audio_content(request, audio_config, content_id, storage_path)
            actual_filename = os.path.basename(generated_content)
            file_path = generated_content
        elif request.contentType == "compiler" and request.contentConfig.get('compiler'):
            compiler_config_dict = request.contentConfig.get('compiler')
            compiler_config = CompilerConfig(**compiler_config_dict)
            generated_content = await generate_compiler_content(request, compiler_config)
            actual_filename = f"{base_filename}.py"
            file_path = os.path.join(storage_path, actual_filename)
            with open(file_path, 'w', encoding='utf-8') as f:
                f.write(generated_content)
        elif request.contentType == "pdf":
            actual_filename = f"{base_filename}.pdf"
            file_path = os.path.join(storage_path, actual_filename)
            await generate_pdf_content_with_path(request, content_id, storage_path, actual_filename)
        elif request.contentType == "ppt":
            actual_filename = f"{base_filename}.pptx"
            file_path = os.path.join(storage_path, actual_filename)
            await generate_ppt_content_with_path(request, content_id, storage_path, actual_filename)
        else:
            raise Exception(f"Unsupported content type: {request.contentType}")
        
        # Get absolute path
        absolute_path = os.path.abspath(file_path)
        storage_directory = os.path.abspath(storage_path)
        
        # Update content status
        update_content_status(content_id, "completed", file_path)
        
        # Update metadata with ACTUAL STORAGE INFO
        content.metadata.update({
            "actualFileName": actual_filename,
            "actualFilePath": absolute_path,
            "storageDirectory": storage_directory
        })
        
        print(f"[CONTENT] Generation completed successfully")
        print(f"[CONTENT] File: {actual_filename}")
        print(f"[CONTENT] Path: {absolute_path}")
        
        return content
        
    except Exception as e:
        print(f"[ERROR] Content generation failed: {str(e)}")
        update_content_status(content_id, "failed", error=str(e))
        raise e

def get_all_content() -> List[GeneratedContent]:
    """Get all content (for debugging/admin purposes)."""
    return list(_content_storage.values())

def track_download(content_id: str, user_id: str, ip_address: str = None) -> None:
    """Track a download event for analytics."""
    if content_id not in _download_tracking:
        _download_tracking[content_id] = []
    
    _download_tracking[content_id].append({
        "timestamp": datetime.now().isoformat(),
        "userId": user_id,
        "ip": ip_address
    })

def get_download_stats(content_id: str) -> Dict:
    """Get download statistics for a content item."""
    downloads = _download_tracking.get(content_id, [])
    return {
        "total_downloads": len(downloads),
        "downloads": downloads,
        "last_downloaded": downloads[-1]["timestamp"] if downloads else None
    }

def clear_content() -> None:
    """Clear all content (for testing purposes)."""
    global _content_storage
    _content_storage = {}