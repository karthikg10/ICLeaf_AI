# backend/app/ppt_generator.py
# PowerPoint generation functions
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from .models import GenerateContentRequest
from . import deps
from .content_utils import (
    content_client, get_rag_context_for_internal_mode,
    validate_rag_context_for_internal_mode, extract_openai_response
)


async def generate_ppt_content(request: GenerateContentRequest, content_id: str, storage_path: str) -> str:
    """Generate PowerPoint with ChatGPT-quality rich content."""
    if not content_client:
        raise Exception("OpenAI client not configured")
    
    # Get RAG context for internal mode
    rag_result = get_rag_context_for_internal_mode(request, top_k=5)
    validate_rag_context_for_internal_mode(rag_result, request)  # Validate context relevance
    rag_context = rag_result.context
    
    ppt_config = request.contentConfig.get('ppt', {})
    num_slides = int(ppt_config.get('num_slides', 10))
    target_audience = ppt_config.get('target_audience', 'general')
    difficulty = ppt_config.get('difficulty', 'medium')
    
    print(f"[PPT] Generating {num_slides} information-rich slides (ChatGPT style)")
    
    system_prompt = f"""You are a world-class presentation expert like ChatGPT. Create {num_slides} RICH, DETAILED slides.

REQUIREMENTS:
- Create EXACTLY {num_slides} slides
- Each slide has a CLEAR TITLE
- Each slide has 4-6 DETAILED bullet points (not just 1-2 words)
- Bullet points are complete sentences with substance
- Include data, examples, explanations
- Professional and engaging
- Suitable for {target_audience} at {difficulty} level

FORMAT FOR EACH SLIDE:
---SLIDE---
TITLE: [Clear, descriptive title]
CONTENT: [4-6 detailed bullet points with substance]
---

SLIDE STRUCTURE:
- Slide 1: Title slide with topic and overview
- Slides 2-{num_slides-1}: Detailed content with examples, data, explanations
- Slide {num_slides}: Summary/Conclusion slide

DETAILED CONTENT RULES:
- Each bullet point should be a complete sentence (20-30 words)
- Include specific examples
- Include relevant statistics or data
- Make it informative, not just pretty
- Rich with actual information"""
    
    if rag_context:
        system_prompt += f"""

IMPORTANT: You are in INTERNAL MODE. Use the provided context from uploaded documents below to create accurate slides. Base your content ONLY on the information provided in the context blocks.

CONTEXT FROM UPLOADED DOCUMENTS:
{rag_context}

Instructions:
- Create slides based on the information in the context above
- Ensure the content is factually accurate to the source material
- Use specific details, examples, and data from the documents
- If the context doesn't cover the topic fully, create slides based on what is available"""

    user_prompt = f"""Create {num_slides} detailed, information-rich slides about: {request.prompt}

Make it like ChatGPT presentations - FULL OF INFORMATION and DETAILED.
Not generic, not simple - RICH CONTENT.

Target audience: {target_audience}
Difficulty: {difficulty}

Each bullet point should have real substance and information.
Each slide should be valuable and detailed.

Generate all {num_slides} slides with rich, detailed content:"""

    try:
        print(f"[PPT] Calling OpenAI API...")
        response = content_client.chat.completions.create(
            model=deps.OPENAI_MODEL,
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_prompt}
            ],
            temperature=0.8,
            max_tokens=5000
        )
        
        content = extract_openai_response(response)
        print(f"[PPT] Got {len(content)} characters")
        
    except Exception as e:
        print(f"[PPT] Error: {e}")
        raise
    
    ppt_path = os.path.join(storage_path, "presentation.pptx")
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    slides_data = []
    current_slide = {"title": "", "bullets": []}
    
    for line in content.split('\n'):
        line = line.strip()
        
        if line == '---SLIDE---':
            continue
        elif line.startswith('TITLE:'):
            current_slide["title"] = line.replace('TITLE:', '').strip()
        elif line.startswith('CONTENT:') or line.startswith('BULLET:') or line.startswith('-'):
            bullet = line.replace('CONTENT:', '').replace('BULLET:', '').replace('-', '', 1).strip()
            if bullet and len(bullet) > 2:
                current_slide["bullets"].append(bullet[:100])
        elif line == '---':
            if current_slide["title"] or current_slide["bullets"]:
                slides_data.append(current_slide)
            current_slide = {"title": "", "bullets": []}
    
    if current_slide["title"] or current_slide["bullets"]:
        slides_data.append(current_slide)
    
    print(f"[PPT] Parsed {len(slides_data)} slides from content")
    
    while len(slides_data) < num_slides:
        slides_data.append({
            "title": f"Additional Information {len(slides_data)}",
            "bullets": ["Detailed content point 1", "Detailed content point 2", "Detailed content point 3"]
        })
    
    blank_layout = prs.slide_layouts[6]
    
    for i, slide_data in enumerate(slides_data[:num_slides]):
        try:
            slide = prs.slides.add_slide(blank_layout)
            
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
            title_frame = title_box.text_frame
            title_frame.word_wrap = True
            title_frame.text = slide_data["title"] or f"Slide {i + 1}"
            
            for paragraph in title_frame.paragraphs:
                paragraph.font.size = Pt(44)
                paragraph.font.bold = True
                paragraph.font.color.rgb = RGBColor(25, 25, 100)
            
            if slide_data["bullets"]:
                content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5.5))
                content_frame = content_box.text_frame
                content_frame.word_wrap = True
                
                for idx, bullet in enumerate(slide_data["bullets"][:6]):
                    if idx == 0:
                        p = content_frame.paragraphs[0]
                    else:
                        p = content_frame.add_paragraph()
                    
                    p.text = f"• {bullet}"
                    p.level = 0
                    p.font.size = Pt(18)
                    p.font.color.rgb = RGBColor(50, 50, 50)
                    p.space_after = Pt(10)
        
        except Exception as e:
            print(f"[PPT] Error on slide {i + 1}: {e}")
            continue
    
    try:
        prs.save(ppt_path)
        print(f"[PPT] Generated {len(prs.slides)} slides: {ppt_path}")
        return ppt_path
    except Exception as e:
        print(f"[PPT] Error saving: {e}")
        raise


async def generate_ppt_content_with_path(
    request: GenerateContentRequest, 
    content_id: str, 
    storage_path: str,
    filename: str
) -> str:
    """Generate PowerPoint with custom filename."""
    if not content_client:
        raise Exception("OpenAI client not configured")
    
    # Get RAG context for internal mode
    rag_result = get_rag_context_for_internal_mode(request, top_k=5)
    validate_rag_context_for_internal_mode(rag_result, request)  # Validate context relevance
    rag_context = rag_result.context
    
    ppt_config = request.contentConfig.get('ppt', {})
    num_slides = int(ppt_config.get('num_slides', 10))
    target_audience = ppt_config.get('target_audience', 'general')
    difficulty = ppt_config.get('difficulty', 'medium')
    
    print(f"[PPT] Generating {num_slides} slides for {target_audience} at {difficulty} level")
    print(f"[PPT] Custom filename: {filename}")
    
    system_prompt = f"""You are a world-class presentation expert like ChatGPT. Create {num_slides} RICH, DETAILED slides.

REQUIREMENTS:
- Create EXACTLY {num_slides} slides
- Each slide has a CLEAR TITLE
- Each slide has 4-6 DETAILED bullet points (not just 1-2 words)
- Bullet points are complete sentences with substance
- Include data, examples, explanations
- Professional and engaging
- Suitable for {target_audience} at {difficulty} level

FORMAT FOR EACH SLIDE:
---SLIDE---
TITLE: [Clear, descriptive title]
CONTENT: [4-6 detailed bullet points with substance]
---

SLIDE STRUCTURE:
- Slide 1: Title slide with topic and overview
- Slides 2-{num_slides-1}: Detailed content with examples, data, explanations
- Slide {num_slides}: Summary/Conclusion slide

DETAILED CONTENT RULES:
- Each bullet point should be a complete sentence (20-30 words)
- Include specific examples
- Include relevant statistics or data
- Make it informative, not just pretty
- Rich with actual information"""
    
    if rag_context:
        system_prompt += f"""

IMPORTANT: You are in INTERNAL MODE. Use the provided context from uploaded documents below to create accurate slides. Base your content ONLY on the information provided in the context blocks.

CONTEXT FROM UPLOADED DOCUMENTS:
{rag_context}

Instructions:
- Create slides based on the information in the context above
- Ensure the content is factually accurate to the source material
- Use specific details, examples, and data from the documents
- If the context doesn't cover the topic fully, create slides based on what is available"""

    user_prompt = f"""Create {num_slides} detailed, information-rich slides about: {request.prompt}

Make it like ChatGPT presentations - FULL OF INFORMATION and DETAILED.
Not generic, not simple - RICH CONTENT.

Target audience: {target_audience}
Difficulty: {difficulty}

Each bullet point should have real substance and information.
Each slide should be valuable and detailed.

Generate all {num_slides} slides with rich, detailed content:"""

    try:
        print(f"[PPT] Calling OpenAI API...")
        response = content_client.chat.completions.create(
            model=deps.OPENAI_MODEL,
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_prompt}
            ],
            temperature=0.8,
            max_tokens=5000
        )
        
        content = extract_openai_response(response)
        print(f"[PPT] Got {len(content)} characters")
        
    except Exception as e:
        print(f"[PPT] Error: {e}")
        raise
    
    ppt_path = os.path.join(storage_path, filename)
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    slides_data = []
    current_slide = {"title": "", "bullets": []}
    
    for line in content.split('\n'):
        line = line.strip()
        
        if line == '---SLIDE---':
            continue
        elif line.startswith('TITLE:'):
            current_slide["title"] = line.replace('TITLE:', '').strip()
        elif line.startswith('CONTENT:') or line.startswith('BULLET:') or line.startswith('-'):
            bullet = line.replace('CONTENT:', '').replace('BULLET:', '').replace('-', '', 1).strip()
            if bullet and len(bullet) > 2:
                current_slide["bullets"].append(bullet[:100])
        elif line == '---':
            if current_slide["title"] or current_slide["bullets"]:
                slides_data.append(current_slide)
            current_slide = {"title": "", "bullets": []}
    
    if current_slide["title"] or current_slide["bullets"]:
        slides_data.append(current_slide)
    
    print(f"[PPT] Parsed {len(slides_data)} slides")
    
    while len(slides_data) < num_slides:
        slides_data.append({
            "title": f"Additional Information {len(slides_data)}",
            "bullets": ["Detailed content point 1", "Detailed content point 2", "Detailed content point 3"]
        })
    
    blank_layout = prs.slide_layouts[6]
    
    for i, slide_data in enumerate(slides_data[:num_slides]):
        try:
            slide = prs.slides.add_slide(blank_layout)
            
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
            title_frame = title_box.text_frame
            title_frame.word_wrap = True
            title_frame.text = slide_data["title"] or f"Slide {i + 1}"
            
            for paragraph in title_frame.paragraphs:
                paragraph.font.size = Pt(44)
                paragraph.font.bold = True
                paragraph.font.color.rgb = RGBColor(25, 25, 100)
            
            if slide_data["bullets"]:
                content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5.5))
                content_frame = content_box.text_frame
                content_frame.word_wrap = True
                
                for idx, bullet in enumerate(slide_data["bullets"][:6]):
                    if idx == 0:
                        p = content_frame.paragraphs[0]
                    else:
                        p = content_frame.add_paragraph()
                    
                    p.text = f"• {bullet}"
                    p.level = 0
                    p.font.size = Pt(18)
                    p.font.color.rgb = RGBColor(50, 50, 50)
                    p.space_after = Pt(10)
        
        except Exception as e:
            print(f"[PPT] Error on slide {i + 1}: {e}")
            continue
    
    try:
        prs.save(ppt_path)
        print(f"[PPT] Generated {len(prs.slides)} slides: {ppt_path}")
        return ppt_path
    except Exception as e:
        print(f"[PPT] Error saving: {e}")
        raise

