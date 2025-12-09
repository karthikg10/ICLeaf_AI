# backend/app/ingest_dir.py
import os
import re
from typing import List, Tuple, Dict

def chunk_text(text: str, max_chars: int = 800, overlap: int = 100) -> List[str]:
    """Chunk text with overlap for RAG. Optimized for presentations."""
    if not text:
        return []
    
    # Remove excessive whitespace but preserve structure
    text = re.sub(r'\n\s*\n', '\n', text)  # Remove blank lines
    text = re.sub(r'[ \t]+', ' ', text)     # Collapse spaces
    text = text.strip()
    
    if not text:
        return []
    
    chunks = []
    n = len(text)
    start = 0
    
    while start < n:
        end = min(start + max_chars, n)
        
        # Try to break at sentence boundary (. or \n) if near end
        if end < n:
            # Look for period or newline within last 100 chars
            last_period = text.rfind('.', max(start, end - 100), end)
            last_newline = text.rfind('\n', max(start, end - 100), end)
            
            break_point = max(last_period, last_newline)
            if break_point > start + 100:  # Only break if we have substantial text before break
                end = break_point + 1  # Include the punctuation
        
        chunk = text[start:end].strip()
        if chunk and len(chunk) > 20:  # Only add meaningful chunks
            chunks.append(chunk)
        
        if end >= n:
            break
        
        start = max(start + 1, end - overlap)
    
    return chunks


def _pretty_title_from_filename(path: str) -> str:
    """Extract clean title from filename."""
    base = os.path.basename(path)
    name, _ = os.path.splitext(base)
    return re.sub(r'[_\-]+', ' ', name).strip()


def _extract_pdf_text_enhanced(path: str) -> str:
    """
    Extract text from PDF with special handling for PowerPoint PDFs and image-based PDFs.
    
    For PPT-as-PDF files and PDFs with images:
    - Uses multiple extraction strategies
    - Preserves structure where possible
    - Uses OCR to extract text from images (even when some text is already extracted)
    - Combines extracted text with OCR text for comprehensive coverage
    """
    import pypdf
    
    all_text = []
    ocr_available = False
    
    # Check if OCR libraries are available
    try:
        import pytesseract
        from pdf2image import convert_from_path
        from PIL import Image
        ocr_available = True
    except ImportError:
        print("[ingest] OCR libraries not available. Install pytesseract, pdf2image, and Pillow for image text extraction.")
    
    try:
        pdf = pypdf.PdfReader(path)
        total_pages = len(pdf.pages)
        
        for page_idx, page in enumerate(pdf.pages):
            # Strategy 1: Standard text extraction
            page_text = page.extract_text() or ""
            extracted_text_length = len(page_text.strip())
            
            # Strategy 2: If extraction is minimal, try extracting from annotations/objects
            if not page_text.strip() or extracted_text_length < 50:
                # For PPT PDFs, text might be in text boxes or special objects
                # Try to extract from page objects if available
                try:
                    if "/Annots" in page:
                        for annot in page["/Annots"]:
                            annot_obj = annot.get_object()
                            if "/Contents" in annot_obj:
                                annot_text = annot_obj["/Contents"]
                                if isinstance(annot_text, str):
                                    page_text += "\n" + annot_text
                                    extracted_text_length = len(page_text.strip())
                except:
                    pass  # Annotation extraction is optional
            
            # Strategy 3: OCR for image-based PDFs and PPTs saved as PDFs
            # Run OCR if:
            # 1. No text was extracted, OR
            # 2. Very little text was extracted (< 100 chars), OR
            # 3. Text seems incomplete (for PPTs saved as PDFs with images)
            should_run_ocr = (
                ocr_available and (
                    not page_text.strip() or 
                    extracted_text_length < 100 or
                    # For PPT PDFs, even if some text exists, images might contain more
                    (extracted_text_length > 0 and extracted_text_length < 200)
                )
            )
            
            ocr_text = ""
            if should_run_ocr:
                try:
                    # Convert page to image and OCR
                    img_list = convert_from_path(
                        path, 
                        first_page=page_idx+1, 
                        last_page=page_idx+1, 
                        dpi=300  # Increased DPI for better OCR accuracy
                    )
                    
                    for img in img_list:
                        if isinstance(img, Image.Image):
                            ocr_result = pytesseract.image_to_string(img, lang='eng')
                            if ocr_result.strip():
                                ocr_text = ocr_result.strip()
                                break
                except Exception as ocr_error:
                    # OCR failed - continue without it
                    print(f"[ingest] OCR failed for page {page_idx + 1}: {ocr_error}")
            
            # Combine extracted text with OCR text
            # Remove duplicates and merge intelligently
            if ocr_text:
                if page_text.strip():
                    # Both extracted text and OCR text exist
                    # Combine them, removing obvious duplicates
                    combined_text = page_text.strip()
                    # Add OCR text if it's significantly different
                    ocr_words = set(ocr_text.lower().split())
                    extracted_words = set(page_text.lower().split())
                    new_words = ocr_words - extracted_words
                    if len(new_words) > 5:  # If OCR found substantial new content
                        combined_text += "\n\n[Text from images]\n" + ocr_text
                    page_text = combined_text
                else:
                    # Only OCR text available
                    page_text = ocr_text
            
            # Clean up and add page text
            if page_text.strip():
                # Add page marker for context
                page_text = page_text.strip()
                all_text.append(f"[Page {page_idx + 1}]\n{page_text}")
        
        return "\n\n".join(all_text)
        
    except Exception as e:
        print(f"[ingest] Error extracting PDF text from {path}: {e}")
        return ""


def _read_file_build_docs(path: str) -> List[Tuple[str, Dict]]:
    """Build document chunks from file with proper metadata."""
    docs: List[Tuple[str, Dict]] = []
    ext = os.path.splitext(path)[1].lower()
    
    try:
        if ext in [".txt", ".md"]:
            with open(path, "r", encoding="utf-8", errors="ignore") as f:
                text = f.read()
            
            title = _pretty_title_from_filename(path)
            
            for i, chunk in enumerate(chunk_text(text)):
                docs.append((chunk, {
                    "id": f"{title}-{i}",
                    "filename": os.path.basename(path),
                    "title": title
                }))
            
            return docs
        
        elif ext == ".pdf":
            title = _pretty_title_from_filename(path)
            
            # Use enhanced extraction for PPT PDFs
            text = _extract_pdf_text_enhanced(path)
            
            if text.strip():
                for i, chunk in enumerate(chunk_text(text)):
                    docs.append((chunk, {
                        "id": f"{title}-{i}",
                        "filename": os.path.basename(path),
                        "title": title
                    }))
            else:
                print(f"[ingest] WARNING: No text extracted from PDF {path}")
            
            return docs
        
        elif ext == ".docx":
            try:
                import docx
                doc = docx.Document(path)
                text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
                title = _pretty_title_from_filename(path)
                
                for i, chunk in enumerate(chunk_text(text)):
                    docs.append((chunk, {
                        "id": f"{title}-{i}",
                        "filename": os.path.basename(path),
                        "title": title
                    }))
            except Exception as e:
                print(f"[ingest] ERROR reading DOCX {path}: {e}")
            
            return docs
        
        elif ext == ".pptx":
            try:
                from pptx import Presentation
                prs = Presentation(path)
                title = _pretty_title_from_filename(path)
                
                slides_text = []
                for slide_num, slide in enumerate(prs.slides):
                    slide_content = []
                    
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            slide_content.append(shape.text)
                    
                    if slide_content:
                        slide_text = f"[Slide {slide_num + 1}]\n" + "\n".join(slide_content)
                        slides_text.append(slide_text)
                
                text = "\n\n".join(slides_text)
                
                for i, chunk in enumerate(chunk_text(text)):
                    docs.append((chunk, {
                        "id": f"{title}-{i}",
                        "filename": os.path.basename(path),
                        "title": title
                    }))
            except Exception as e:
                print(f"[ingest] ERROR reading PPTX {path}: {e}")
            
            return docs
        
        else:
            print(f"[ingest] Unsupported file type: {ext}")
            return docs
    
    except Exception as e:
        print(f"[ingest] ERROR processing file {path}: {e}")
        return docs


def build_docs_for_dir(root_dir: str) -> List[Tuple[str, Dict]]:
    """Build all documents from directory."""
    out: List[Tuple[str, Dict]] = []
    
    if not root_dir or not os.path.isdir(root_dir):
        print(f"[ingest] Directory not found: {root_dir}")
        return out
    
    print(f"[ingest] Scanning {root_dir} ...")
    
    for dirpath, _, filenames in os.walk(root_dir):
        for fn in filenames:
            path = os.path.join(dirpath, fn)
            ext = os.path.splitext(path)[1].lower()
            
            if ext not in [".txt", ".md", ".pdf", ".docx", ".pptx"]:
                continue
            
            print(f"[ingest] Processing: {os.path.basename(path)}")
            
            docs = _read_file_build_docs(path)
            
            if docs:
                print(f"[ingest] OK {os.path.basename(path)}: +{len(docs)} chunks")
                out.extend(docs)
            else:
                print(f"[ingest] WARNING: No chunks extracted from {path}")
    
    print(f"[ingest] Total chunks from {root_dir}: {len(out)}")
    return out


# Import RAG store for adding documents
import app.rag_store_chromadb as rag


def ingest_dir(root_dir: str, subject_id: str = None, topic_id: str = None, uploaded_by: str = None) -> int:
    """Ingest all documents from directory into RAG store."""
    docs = build_docs_for_dir(root_dir)
    
    if not docs:
        print(f"[ingest] No documents to ingest from {root_dir}")
        return 0
    
    # Add metadata if provided
    if subject_id or topic_id or uploaded_by:
        enhanced_docs = []
        for text, meta in docs:
            enhanced_meta = meta.copy()
            if subject_id:
                enhanced_meta["subjectId"] = subject_id
            if topic_id:
                enhanced_meta["topicId"] = topic_id
            if uploaded_by:
                enhanced_meta["uploadedBy"] = uploaded_by
            enhanced_docs.append((text, enhanced_meta))
        docs = enhanced_docs
    
    # Add to RAG store
    rag.add_documents(docs)
    print(f"[ingest] Added {len(docs)} chunks to RAG store")
    
    return len(docs)
